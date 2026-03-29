"""Extract visual DNA (fonts, colors, logo, footer) from a template PPTX."""
from __future__ import annotations

import re
from collections import Counter, defaultdict
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_template_transfer.models import TemplateStyle
from pptx_template_transfer.helpers import FOOTER_PATTERNS, NSMAP


def _extract_theme_fonts(prs: Presentation) -> tuple[str, str]:
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    heading = "Montserrat"
    body = "Lato"
    try:
        master = prs.slide_masters[0]
        theme_el = master.element.find(f".//{{{ns_a}}}theme")
        if theme_el is None:
            for rel in master.part.rels.values():
                if "theme" in str(rel.reltype).lower():
                    theme_xml = rel.target_part.blob
                    theme_el = etree.fromstring(theme_xml)
                    break
        if theme_el is not None:
            major = theme_el.find(f".//{{{ns_a}}}majorFont")
            minor = theme_el.find(f".//{{{ns_a}}}minorFont")
            if major is not None:
                lat = major.find(f"{{{ns_a}}}latin")
                if lat is not None and lat.get("typeface"):
                    heading = lat.get("typeface")
            if minor is not None:
                lat = minor.find(f"{{{ns_a}}}latin")
                if lat is not None and lat.get("typeface"):
                    body = lat.get("typeface")
    except Exception:
        pass

    if heading == body:
        large_fonts: Counter[str] = Counter()
        body_fonts: Counter[str] = Counter()
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            if run.font.size and run.font.size.pt >= 18:
                                large_fonts[run.font.name] += len(run.text)
                            else:
                                body_fonts[run.font.name] += len(run.text)
        if large_fonts:
            heading = large_fonts.most_common(1)[0][0]
        if body_fonts:
            body = body_fonts.most_common(1)[0][0]
    return heading, body


def _extract_colors(prs: Presentation) -> dict[str, str]:
    color_freq: Counter[str] = Counter()
    bg_color = "F7F8FB"

    try:
        fill = prs.slides[0].background.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc.type is not None and fc.rgb:
                bg_color = str(fc.rgb)
    except Exception:
        pass

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        c = run.font.color
                        if c and c.type is not None and c.rgb:
                            color_freq[str(c.rgb)] += len(run.text)
                    except (AttributeError, TypeError):
                        pass

    text_color = "111827"
    muted_color = "475569"
    primary_color = "2563EB"
    secondary_color = "F97316"

    dark_colors = []
    saturated_accents = []
    muted_accents = []
    for c, freq in color_freq.most_common(20):
        r, g, b = int(c[:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        brightness = (r + g + b) / 3
        max_ch, min_ch = max(r, g, b), min(r, g, b)
        saturation = (max_ch - min_ch) / max_ch if max_ch > 0 else 0

        if brightness < 80:
            dark_colors.append((c, freq))
        elif c != bg_color and brightness < 240:
            if saturation > 0.4:
                saturated_accents.append((c, freq))
            elif brightness < 160:
                muted_accents.append((c, freq))

    if dark_colors:
        text_color = dark_colors[0][0]
    if muted_accents:
        muted_color = muted_accents[0][0]
    elif len(dark_colors) >= 2:
        muted_color = dark_colors[1][0]
    if saturated_accents:
        primary_color = saturated_accents[0][0]
        if len(saturated_accents) >= 2:
            secondary_color = saturated_accents[1][0]

    return {
        "text": text_color, "muted": muted_color,
        "primary": primary_color, "secondary": secondary_color,
        "background": bg_color, "card": "FFFFFF", "line": "D1D5DB",
    }


def _extract_logo(prs: Presentation) -> tuple[bytes | None, str, int, int]:
    img_map: dict[int, list] = defaultdict(list)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13:
                    blob = shape.image.blob
                    key = hash(blob[:200])
                    img_map[key].append((blob, shape.image.content_type,
                                         shape.width, shape.height))
            except Exception:
                pass

    best_blob = None
    best_ct = "image/png"
    best_w = best_h = 0
    best_count = 0
    for key, occurrences in img_map.items():
        if len(occurrences) > best_count:
            best_count = len(occurrences)
            b, ct, w, h = occurrences[0]
            best_blob, best_ct, best_w, best_h = b, ct, w, h

    if best_count < 2:
        return None, "image/png", 0, 0
    return best_blob, best_ct, best_w, best_h


def _extract_footer_text(prs: Presentation) -> str:
    """Extract the recurring footer/company text from the template.

    Only text that appears in the bottom 10% of **multiple** slides is treated
    as a true footer line.  A sentence that appears on just one slide is
    ordinary content that happens to sit low on the page — not a footer.
    """
    footer_texts: Counter[str] = Counter()
    sh = prs.slide_height
    for slide in prs.slides:
        seen_on_slide: set[str] = set()
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            bottom = (shape.top or 0) + (shape.height or 0)
            if bottom / sh > 0.90:
                t = shape.text_frame.text.strip()
                if (
                    t
                    and not FOOTER_PATTERNS.match(t)
                    and not re.match(r"^Page\s*\d+$", t, re.I)
                    and t not in seen_on_slide
                ):
                    seen_on_slide.add(t)
                    footer_texts[t] += 1

    # Require the text to appear on at least 2 slides (true footer vs one-off
    # body content that sits in the bottom zone of a single slide).
    if footer_texts:
        best, count = footer_texts.most_common(1)[0]
        if count >= 2:
            return best
    return ""


def analyze_template(template_path: Path) -> TemplateStyle:
    """Analyze a template PPTX and extract its visual DNA."""
    prs = Presentation(str(template_path))
    style = TemplateStyle()
    style.slide_width = prs.slide_width
    style.slide_height = prs.slide_height

    style.heading_font, style.body_font = _extract_theme_fonts(prs)

    colors = _extract_colors(prs)
    style.color_primary = colors["primary"]
    style.color_secondary = colors["secondary"]
    style.color_text = colors["text"]
    style.color_muted = colors["muted"]
    style.color_background = colors["background"]
    style.color_card = colors["card"]
    style.color_line = colors["line"]

    blob, ct, w, h = _extract_logo(prs)
    style.logo_blob = blob
    style.logo_content_type = ct
    style.logo_width = w
    style.logo_height = h

    style.footer_company = _extract_footer_text(prs)
    style.footer_has_confidential = True
    style.footer_has_page_number = True

    # Mine layout patterns (Phase 2)
    from pptx_template_transfer.analysis.layout_patterns import mine_layout_patterns
    style.patterns = mine_layout_patterns(prs)

    return style
