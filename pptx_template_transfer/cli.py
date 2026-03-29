"""CLI entry point and public API for PPTX Template Transfer."""
from __future__ import annotations

import argparse
import io
import json
import logging
import sys
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu

from pptx_template_transfer.models import TransferConfig, Thresholds
from pptx_template_transfer.helpers import text_of, word_count, max_font_pt, shape_area_pct
from pptx_template_transfer.analysis.slide_classifier import (
    classify_all_shapes, get_slide_zones,
)
from pptx_template_transfer.extraction.content_extractor import extract_content
from pptx_template_transfer.transform.slide_builder import apply_recreate
from pptx_template_transfer.transform.clone_injector import apply_design


log = logging.getLogger("pptx_template_transfer")


# ============================================================================
# AUTO-DETECTION
# ============================================================================

def detect_mode(template_path: Path) -> str:
    return "recreate"


# ============================================================================
# INPUT VALIDATION
# ============================================================================

def _validate_input(path: Path, label: str) -> None:
    if not path.exists():
        print(f"Error: {label} not found: {path}", file=sys.stderr)
        sys.exit(1)
    if not path.suffix.lower() == ".pptx":
        print(f"Error: {label} must be a .pptx file: {path}", file=sys.stderr)
        sys.exit(1)
    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            if "[Content_Types].xml" not in zf.namelist():
                print(f"Error: {label} is not a valid PPTX: {path}", file=sys.stderr)
                sys.exit(1)
    except zipfile.BadZipFile:
        print(f"Error: {label} is corrupt or not a ZIP archive: {path}", file=sys.stderr)
        sys.exit(1)

    try:
        prs = Presentation(str(path))
        if len(prs.slides) == 0:
            print(f"Warning: {label} has 0 slides: {path}", file=sys.stderr)
    except Exception as exc:
        print(f"Error: Cannot load {label}: {exc}", file=sys.stderr)
        sys.exit(1)


# ============================================================================
# LAYOUT MODE
# ============================================================================

def apply_layout(
    template_path: Path, content_path: Path, output_path: Path,
    config: TransferConfig,
) -> dict[str, Any]:
    print("\n[layout] Loading presentations...")
    template_prs = Presentation(str(template_path))
    print(f"  Template layouts: {[l.name for l in template_prs.slide_layouts]}")
    print("  [layout] Using design-mode pipeline (python-pptx layout limitation)")
    return apply_design(template_path, content_path, output_path, config)


# ============================================================================
# PUBLIC API
# ============================================================================

def transfer(
    template: Path, content: Path,
    output: Path | None = None,
    config: TransferConfig | None = None,
) -> dict[str, Any]:
    """Programmatic API for template transfer. Returns the report dict."""
    if config is None:
        config = TransferConfig()

    mode = config.mode or detect_mode(template)
    if output is None:
        output = Path("output.pptx")

    if mode == "recreate":
        return apply_recreate(template, content, output, config)
    if mode == "design" or mode == "clone":
        return apply_design(template, content, output, config)
    return apply_layout(template, content, output, config)


# ============================================================================
# CLI: ANALYSIS MODES
# ============================================================================

def _cli_analyze(pptx_path: Path) -> None:
    prs = Presentation(str(pptx_path))
    sw, sh = prs.slide_width, prs.slide_height
    th = Thresholds()

    print(f"\nAnalyzing: {pptx_path}")
    print(f"Slides: {len(prs.slides)}, Size: {Emu(sw).inches:.1f}\" x {Emu(sh).inches:.1f}\"\n")

    for i, slide in enumerate(prs.slides):
        zones = get_slide_zones(slide, sw, sh, th)
        classifications = classify_all_shapes(slide, sw, sh, th)
        total = len(list(slide.shapes))

        print(f"Slide {i+1}/{len(prs.slides)}: {total} shapes")
        print(f"  Zones: title={len(zones['title'])}, body={len(zones['body'])}, "
              f"footer={len(zones['footer'])}, decorative={len(zones['decorative'])}")

        for shape, role, conf in classifications:
            text = text_of(shape)
            preview = f' "{text[:50]}"' if text else ""
            area = round(shape_area_pct(shape, sw, sh), 1)
            top = round((shape.top or 0) / sh * 100 if sh else 0)
            print(f'    [{role:11s} {conf:.2f}] "{shape.name}" '
                  f'({area}% area, top {top}%){preview}')
        print()


def _cli_extract(pptx_path: Path) -> None:
    prs = Presentation(str(pptx_path))
    sw, sh = prs.slide_width, prs.slide_height
    th = Thresholds()
    ct = len(prs.slides)

    result = []
    for i, slide in enumerate(prs.slides):
        cd = extract_content(slide, i, ct, sw, sh, th)
        slide_data = {
            "slide": i + 1,
            "slide_type": cd.slide_type,
            "title": cd.title,
            "word_count": cd.word_count,
            "body_paragraphs": [
                {"text": p.text, "bold": p.bold, "level": p.level}
                for p in cd.body_paragraphs
            ],
            "tables": [t["data"] for t in cd.tables],
            "images": [
                {"width": img[1], "height": img[2], "blob_size": len(img[0])}
                for img in cd.images
            ],
            "has_chart": cd.has_chart,
            "notes": cd.notes if cd.notes else None,
        }
        result.append(slide_data)

    print(json.dumps(result, indent=2, ensure_ascii=False))


# ============================================================================
# CLI: MAIN
# ============================================================================

def main() -> None:
    parser = argparse.ArgumentParser(
        description="PPTX Template Transfer - apply one deck's visual design to another's content.",
    )
    parser.add_argument("template_pptx", type=Path,
                        help="Template/target PPTX (style source). Also: --target alias.")
    parser.add_argument("content_pptx", type=Path, nargs="?", default=None,
                        help="Content/source PPTX (content source). Also: --source alias.")
    parser.add_argument("output_pptx", type=Path, nargs="?", default=None,
                        help="Output PPTX path")
    parser.add_argument("--source", type=Path, default=None,
                        help="Alias for content_pptx (source deck)")
    parser.add_argument("--target", type=Path, default=None,
                        help="Alias for template_pptx (target style deck)")
    parser.add_argument("--output", type=Path, default=None,
                        help="Alias for output_pptx")
    parser.add_argument("--mode", choices=["recreate", "design", "clone", "layout"], default=None)
    parser.add_argument("--slide-map", type=Path, default=None,
                        help='JSON: {"1": 3, "2": 1, ...}')
    parser.add_argument("--layout-map", type=Path, default=None)
    parser.add_argument("--verbose", "-v", action="store_true")
    parser.add_argument("--report", type=Path, default=None,
                        help="Write JSON diagnostics report to this path")
    parser.add_argument("--quality-report", type=Path, default=None,
                        help="Write quality validation report to this path")
    parser.add_argument("--no-notes", action="store_true",
                        help="Skip speaker notes transfer")
    parser.add_argument("--analyze", action="store_true",
                        help="Analyze a single PPTX: classify every shape on every slide")
    parser.add_argument("--extract", action="store_true",
                        help="Extract structured content from a single PPTX as JSON")

    args = parser.parse_args()

    if hasattr(sys.stdout, "buffer"):
        sys.stdout = io.TextIOWrapper(
            sys.stdout.buffer, encoding="utf-8", errors="replace",
        )

    level = logging.DEBUG if args.verbose else logging.INFO
    logging.basicConfig(level=level, format="%(message)s", stream=sys.stdout)

    # --- Handle --source/--target aliases ---
    template_path = args.target or args.template_pptx
    content_path = args.source or args.content_pptx
    output_path = args.output or args.output_pptx

    # --- Single-file analysis modes ---
    if args.analyze:
        _validate_input(template_path, "Input")
        _cli_analyze(template_path)
        return

    if args.extract:
        _validate_input(template_path, "Input")
        _cli_extract(template_path)
        return

    # --- Transfer mode ---
    if not content_path or not output_path:
        parser.error("content_pptx and output_pptx are required for transfer mode")

    _validate_input(template_path, "Template")
    _validate_input(content_path, "Content")

    slide_map = None
    if args.slide_map and args.slide_map.exists():
        slide_map = json.loads(args.slide_map.read_text())

    config = TransferConfig(
        mode=args.mode, verbose=args.verbose, slide_map=slide_map,
        preserve_notes=not args.no_notes, report_path=args.report,
    )

    mode = config.mode or detect_mode(template_path)
    if config.mode is None:
        print(f"Auto-detected mode: {mode}")
    config = TransferConfig(
        mode=mode, verbose=config.verbose, slide_map=config.slide_map,
        preserve_notes=config.preserve_notes, thresholds=config.thresholds,
        report_path=config.report_path,
    )

    if mode == "recreate":
        report = apply_recreate(template_path, content_path, output_path, config)
    elif mode in ("design", "clone"):
        report = apply_design(template_path, content_path, output_path, config)
    else:
        report = apply_layout(template_path, content_path, output_path, config)

    if args.report:
        clean_report = json.loads(json.dumps(report, default=str))
        args.report.write_text(json.dumps(clean_report, indent=2))
        print(f"Report written to {args.report}")

    # Quality report
    if args.quality_report:
        from pptx_template_transfer.validation.quality_report import generate_quality_report
        from pptx_template_transfer.extraction.content_extractor import extract_all_content
        output_prs = Presentation(str(output_path))
        content_list = extract_all_content(content_path, config.thresholds)
        qr = generate_quality_report(output_prs, content_list, config)
        qr_dict = {
            "overall_score": qr.overall_score,
            "native_count": qr.native_count,
            "fallback_count": qr.fallback_count,
            "warnings": qr.warnings,
            "slides": [
                {
                    "index": s.slide_index + 1,
                    "build_method": s.build_method,
                    "content_coverage_pct": s.content_coverage_pct,
                    "font_warnings": s.font_warnings,
                    "overlap_count": len(s.overlap_issues),
                    "bounds_issues": len(s.bounds_issues),
                    "needs_review": s.needs_manual_review,
                    "review_reasons": s.review_reasons,
                }
                for s in qr.slides
            ],
        }
        args.quality_report.write_text(json.dumps(qr_dict, indent=2))
        print(f"Quality report written to {args.quality_report}")
        print(f"  Overall score: {qr.overall_score:.0f}/100")
        if qr.warnings:
            for w in qr.warnings[:10]:
                print(f"  {w}")


if __name__ == "__main__":
    main()
