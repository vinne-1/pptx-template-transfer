"""Recreate-mode slide builder — constructs slides from scratch.

Provides all functions needed to build output slides in recreate mode:
background, decorations, header/footer, text, tables, images, and the
main orchestrator that drives the full pipeline.
"""
from __future__ import annotations

import io
import json
import logging
import traceback
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from pptx_template_transfer.helpers import NSMAP, rgb, style_runs
from pptx_template_transfer.models import (
    BrandingPolicy,
    ContentData,
    ParagraphData,
    TemplateStyle,
    TextBlock,
    TransferConfig,
)

# Module-level branding context (set by build_slide / apply_recreate)
_current_branding: BrandingPolicy = BrandingPolicy()

# ---------------------------------------------------------------------------
# Section label generation
# ---------------------------------------------------------------------------

# Map slide types to clean, professional section labels
_TYPE_LABELS: dict[str, str] = {
    "title": "",
    "section": "",
    "agenda": "AGENDA",
    "closing": "SUMMARY",
    "data": "DATA OVERVIEW",
    "metrics_dashboard": "KEY METRICS",
    "comparison": "ANALYSIS",
    "image_heavy": "EVIDENCE",
    "content_narrative": "",  # derived from title
    "content_bullets": "",
    "process_flow": "PROCESS",
    "timeline": "TIMELINE",
    "toc": "TABLE OF CONTENTS",
}


def _generate_section_label(content: ContentData) -> str:
    """Generate a clean, professional section label from source content.

    Rules:
    - Use type-based label if available
    - Otherwise derive from title: take meaningful words, max 3 words
    - Strip punctuation artifacts (?, -, trailing prepositions)
    - Omit label if confidence is low (result too short or weird)
    """
    # Incident heuristic overrides type-based label
    if _is_incident_slide(content):
        return "INCIDENT"

    # Check type-based label first
    type_label = _TYPE_LABELS.get(content.slide_type, "")
    if type_label:
        return type_label

    if not content.title:
        return "OVERVIEW"

    title = content.title.strip()

    # Special case: if title contains a dash separator, use the part before it
    # e.g., "Incident Overview – SSScheduler.exe" -> "INCIDENT OVERVIEW"
    for sep in [" – ", " - ", " — ", ": "]:
        if sep in title:
            title = title.split(sep)[0].strip()
            break

    # Take up to 3 words, stripping articles and prepositions from the end
    words = title.split()
    _TRAILING_JUNK = {"the", "a", "an", "of", "for", "in", "on", "to", "and", "&", "with", "by"}
    label_words: list[str] = []
    for w in words:
        if len(label_words) >= 3:
            break
        label_words.append(w)

    # Remove trailing junk words
    while label_words and label_words[-1].lower() in _TRAILING_JUNK:
        label_words.pop()

    if not label_words:
        return ""

    label = " ".join(label_words).upper()

    # Strip trailing punctuation artifacts
    label = label.rstrip("?!:;,.-–—")

    # Quality gate: label must be 2+ chars and not just a number
    if len(label) < 2 or label.isdigit():
        return ""

    return label

log = logging.getLogger("pptx_template_transfer")

# ---------------------------------------------------------------------------
# Slide‑level decoration helpers (private)
# ---------------------------------------------------------------------------


def _add_background(slide, style: TemplateStyle) -> None:
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(style.color_background)


def _add_decorative_shapes(slide, style: TemplateStyle) -> None:
    """Add corner decorative shapes matching the template style."""
    sw, sh = style.slide_width, style.slide_height

    try:
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE

    # Bottom-right ellipse (large, subtle)
    ellipse_size = int(sw * 0.16)
    ellipse = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        sw - int(ellipse_size * 0.7),
        sh - int(ellipse_size * 0.7),
        ellipse_size, ellipse_size,
    )
    fill = ellipse.fill
    fill.solid()
    fill.fore_color.rgb = rgb(style.color_primary)
    # Set transparency via XML
    ns_a = NSMAP["a"]
    solid_fill = ellipse._element.find(f".//{{{ns_a}}}solidFill")
    if solid_fill is not None:
        color_el = solid_fill[0] if len(solid_fill) else None
        if color_el is not None:
            alpha = etree.SubElement(color_el, f"{{{ns_a}}}alpha")
            alpha.set("val", "25000")  # 25% opacity
    ellipse.line.fill.background()  # no border

    # Top-right triangle
    tri_w = int(sw * 0.17)
    tri_h = int(sh * 0.33)
    triangle = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        sw - tri_w, 0,
        tri_w, tri_h,
    )
    fill = triangle.fill
    fill.solid()
    fill.fore_color.rgb = rgb(style.color_line)
    solid_fill = triangle._element.find(f".//{{{ns_a}}}solidFill")
    if solid_fill is not None:
        color_el = solid_fill[0] if len(solid_fill) else None
        if color_el is not None:
            alpha = etree.SubElement(color_el, f"{{{ns_a}}}alpha")
            alpha.set("val", "20000")  # 20% opacity
    triangle.line.fill.background()
    # Flip horizontal so hypotenuse faces left
    triangle.rotation = 180.0


def _add_header(
    slide, style: TemplateStyle, section_label: str,
) -> None:
    """Add accent line and section label above the title area."""
    sw = style.slide_width
    left = int(sw * 0.054)
    # Blue accent line
    line_w = int(sw * 0.038)
    line_h = Pt(3)
    line_top = int(style.slide_height * 0.075)
    line_shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, line_top, line_w, int(line_h),
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = rgb(style.color_primary)
    line_shape.line.fill.background()

    # Section label (ALL-CAPS, small, primary color)
    if section_label:
        # Sanitize non-ASCII for broad viewer compat
        safe_label = section_label.upper().encode("ascii", "replace").decode("ascii")
        lbl = slide.shapes.add_textbox(
            left, line_top - Pt(16), int(sw * 0.6), Pt(14),
        )
        tf = lbl.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = safe_label
        style_runs(p, font_name=style.heading_font, font_size_pt=8,
                    bold=True, color_hex=style.color_primary)


def _add_logo(slide, style: TemplateStyle) -> None:
    """Add logo image in the top-left area."""
    if not style.logo_blob:
        return
    try:
        left = int(style.slide_width * 0.024)
        top = int(style.slide_height * 0.030)
        # Scale logo to reasonable size
        max_w = int(style.slide_width * 0.12)
        max_h = int(style.slide_height * 0.05)
        w, h = style.logo_width, style.logo_height
        if w > 0 and h > 0:
            scale = min(max_w / w, max_h / h, 1.0)
            w, h = int(w * scale), int(h * scale)
        else:
            w, h = max_w, max_h
        slide.shapes.add_picture(io.BytesIO(style.logo_blob), left, top, w, h)
    except Exception as exc:
        log.warning("Logo placement failed: %s", exc)


def _add_footer(
    slide, style: TemplateStyle,
    slide_number: int, total_slides: int,
    branding: BrandingPolicy | None = None,
) -> None:
    """Add footer bar with company name, confidential, page number.

    Respects BrandingPolicy:
    - mode="target" (default): use target template footer text
    - footer_company_override: explicit company name
    - confidentiality_label: custom label (default "Confidential")
    """
    sw, sh = style.slide_width, style.slide_height
    footer_top = int(sh * 0.94)
    branding = branding or BrandingPolicy()

    # Resolve company text
    company_text = branding.footer_company_override or style.footer_company

    # Company name (left)
    if company_text:
        tb = slide.shapes.add_textbox(
            int(sw * 0.04), footer_top, int(sw * 0.35), Pt(12),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = company_text
        style_runs(p, font_name=style.body_font, font_size_pt=7,
                    color_hex=style.color_muted)

    # Confidential (center-right)
    if style.footer_has_confidential:
        conf_label = branding.confidentiality_label or "Confidential"
        tb = slide.shapes.add_textbox(
            int(sw * 0.42), footer_top, int(sw * 0.2), Pt(12),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = conf_label
        style_runs(p, font_name=style.body_font, font_size_pt=7,
                    color_hex=style.color_muted)

    # Page number (right)
    if style.footer_has_page_number:
        tb = slide.shapes.add_textbox(
            int(sw * 0.90), footer_top, int(sw * 0.07), Pt(12),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"Page {slide_number:02d}"
        p.alignment = PP_ALIGN.RIGHT
        style_runs(p, font_name=style.body_font, font_size_pt=7,
                    color_hex=style.color_muted)


# ---------------------------------------------------------------------------
# Content placement helpers (private)
# ---------------------------------------------------------------------------


def _add_title_text(
    slide, style: TemplateStyle, title: str,
    left: int, top: int, width: int,
    font_size_pt: float = 22.0, bold: bool = True,
) -> None:
    """Add title text box."""
    tb = slide.shapes.add_textbox(left, top, width, Pt(font_size_pt + 8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    style_runs(p, font_name=style.heading_font, font_size_pt=font_size_pt,
                bold=bold, color_hex=style.color_text)


def _add_body_text(
    slide, style: TemplateStyle,
    paragraphs: list[ParagraphData],
    left: int, top: int, width: int, max_height: int,
) -> None:
    """Add body text with density-aware font scaling and section grouping."""
    if not paragraphs:
        return

    # --- Density-aware font scaling ---
    # Estimate total lines (accounting for wrapped text at ~80 chars per line)
    total_lines = 0
    for pd in paragraphs:
        text_len = len(pd.text) if pd.text else 0
        total_lines += max(1, text_len // 80 + 1)
    # Add spacing for bold headings
    total_lines += sum(1 for pd in paragraphs if pd.bold)

    # Scale factor: base sizes for <=12 lines, shrink for denser content
    if total_lines <= 12:
        scale = 1.0
    elif total_lines <= 20:
        scale = 0.85
    elif total_lines <= 30:
        scale = 0.72
    else:
        scale = 0.62

    heading_pt = max(10.0, round(14.0 * scale, 1))
    body_pt = max(8.0, round(12.0 * scale, 1))
    sub_pt = max(7.0, round(11.0 * scale, 1))
    spacing_before_heading = max(Pt(4), Pt(round(10 * scale)))

    tb = slide.shapes.add_textbox(left, top, width, max_height)
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for pd in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.text = pd.text
        if pd.bold:
            p.space_before = spacing_before_heading
            style_runs(p, font_name=style.heading_font, font_size_pt=heading_pt,
                        bold=True, color_hex=style.color_text)
        elif pd.level > 0:
            p.level = pd.level
            style_runs(p, font_name=style.body_font, font_size_pt=sub_pt,
                        color_hex=style.color_muted)
        else:
            style_runs(p, font_name=style.body_font, font_size_pt=body_pt,
                        color_hex=style.color_text)


def _add_text_blocks(
    slide, style: TemplateStyle,
    text_blocks: list[TextBlock],
) -> None:
    """Recreate positioned text blocks preserving spatial layout from source."""
    sw, sh = style.slide_width, style.slide_height

    for block in text_blocks:
        left = int(sw * block.left_pct / 100)
        top = int(sh * block.top_pct / 100)
        width = int(sw * block.width_pct / 100)
        height = int(sh * block.height_pct / 100)

        # Ensure minimum dimensions
        width = max(width, int(sw * 0.05))
        height = max(height, int(sh * 0.02))

        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True

        first = True
        for pd in block.paragraphs:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            p.text = pd.text

            if block.is_label:
                # Small labels: numbers, tags — use accent color, compact
                style_runs(p, font_name=style.body_font, font_size_pt=10,
                            bold=True, color_hex=style.color_primary)
            elif block.is_heading or pd.bold:
                # Section headings
                p.space_before = Pt(4)
                style_runs(p, font_name=style.heading_font, font_size_pt=13,
                            bold=True, color_hex=style.color_text)
            elif pd.level > 0:
                p.level = pd.level
                style_runs(p, font_name=style.body_font, font_size_pt=10,
                            color_hex=style.color_muted)
            else:
                style_runs(p, font_name=style.body_font, font_size_pt=11,
                            color_hex=style.color_text)


def _add_table(
    slide, style: TemplateStyle,
    table_data: list[list[str]],
    left: int, top: int, width: int, max_height: int,
) -> None:
    """Build a styled table from content data."""
    if not table_data or not table_data[0]:
        return
    rows, cols = len(table_data), len(table_data[0])
    row_height = min(Pt(24), max_height // rows) if rows else Pt(24)

    shape = slide.shapes.add_table(rows, cols, left, top, width, rows * row_height)
    table = shape.table

    # Style each cell
    for ri in range(rows):
        for ci in range(min(cols, len(table_data[ri]) if ri < len(table_data) else 0)):
            cell = table.cell(ri, ci)
            cell.text = (
                table_data[ri][ci]
                if ri < len(table_data) and ci < len(table_data[ri])
                else ""
            )

            # Format
            for p in cell.text_frame.paragraphs:
                style_runs(p, font_name=style.body_font, font_size_pt=9,
                            color_hex=style.color_text)

            # Header row styling
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(style.color_primary)
                for p in cell.text_frame.paragraphs:
                    style_runs(p, font_name=style.body_font, font_size_pt=9,
                                bold=True, color_hex="FFFFFF")
            elif ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(style.color_background)


def _add_content_images(
    slide, style: TemplateStyle,
    images: list[tuple],
    start_top: int,
    has_body_text: bool = True,
    preserve_position: bool = False,
) -> None:
    """Place content images on the slide with adaptive layout."""
    sw, sh = style.slide_width, style.slide_height

    # Separate positioned vs fallback images
    positioned: list[tuple] = []
    fallback: list[tuple] = []
    for img_tuple in images:
        blob = img_tuple[0]
        orig_w, orig_h = img_tuple[1], img_tuple[2]
        orig_left = img_tuple[3] if len(img_tuple) > 3 else 0
        orig_top = img_tuple[4] if len(img_tuple) > 4 else start_top

        if preserve_position and orig_left > 0 and orig_top > 0:
            positioned.append((blob, orig_w, orig_h, orig_left, orig_top))
        else:
            fallback.append((blob, orig_w, orig_h))

    # Place positioned images at original coordinates
    for blob, ow, oh, ol, ot in positioned:
        try:
            slide.shapes.add_picture(io.BytesIO(blob), ol, ot, ow, oh)
        except Exception:
            pass

    if not fallback:
        return

    # Adaptive layout for fallback images based on count
    n = len(fallback)
    margin_right = int(sw * 0.04)
    avail_h = sh - start_top - int(sh * 0.08)

    if n <= 2:
        # Stack vertically on the right
        max_w_frac = 0.42 if has_body_text else 0.70
        max_w = int(sw * max_w_frac)
        per_img_h = avail_h // n - Pt(4)
        for blob, orig_w, orig_h in fallback:
            max_h = per_img_h
            w, h = orig_w, orig_h
            if w > 0 and h > 0:
                scale = min(max_w / w, max_h / h, 1.0)
                w, h = int(w * scale), int(h * scale)
            else:
                w, h = max_w, max_h
            try:
                left = sw - w - margin_right
                slide.shapes.add_picture(io.BytesIO(blob), left, start_top, w, h)
                start_top += h + Pt(6)
            except Exception:
                pass
    else:
        # Grid layout for 3+ images
        cols = 2 if n <= 4 else 3
        rows_count = (n + cols - 1) // cols
        grid_w = int(sw * (0.42 if has_body_text else 0.85))
        grid_left = sw - grid_w - margin_right
        cell_w = (grid_w - Pt(4) * (cols - 1)) // cols
        cell_h = min((avail_h - Pt(4) * (rows_count - 1)) // rows_count, cell_w)

        for idx, (blob, orig_w, orig_h) in enumerate(fallback):
            r, c = divmod(idx, cols)
            img_left = grid_left + c * (cell_w + Pt(4))
            img_top = start_top + r * (cell_h + Pt(4))
            w, h = orig_w, orig_h
            if w > 0 and h > 0:
                scale = min(cell_w / w, cell_h / h, 1.0)
                w, h = int(w * scale), int(h * scale)
            else:
                w, h = cell_w, cell_h
            try:
                slide.shapes.add_picture(io.BytesIO(blob), img_left, img_top, w, h)
            except Exception:
                pass


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------


def _find_blank_layout(prs: Presentation):
    """Find the blank slide layout (no placeholders, or named 'Blank')."""
    # Prefer layout named "Blank" or "blank"
    for layout in prs.slide_layouts:
        if layout.name.strip().lower() in ("blank", "empty"):
            return layout
    # Fallback: layout with fewest placeholders
    best = prs.slide_layouts[0]
    best_count = len(best.placeholders)
    for layout in prs.slide_layouts:
        if len(layout.placeholders) < best_count:
            best = layout
            best_count = len(layout.placeholders)
    return best


# ---------------------------------------------------------------------------
# Slide‑type builders (private)
# ---------------------------------------------------------------------------


def _build_title_slide(
    slide, style: TemplateStyle, content: ContentData,
    slide_number: int, total_slides: int,
) -> None:
    """Build a title/cover slide."""
    sw, sh = style.slide_width, style.slide_height
    # NOTE: logo already added by build_slide() — no duplicate call here

    # Use positioned text_blocks when available (preserves original layout)
    if content.text_blocks:
        _add_text_blocks(slide, style, content.text_blocks)
    else:
        # Fallback: centered title + subtitle
        title_width = int(sw * 0.7)
        title_left = (sw - title_width) // 2
        title_top = int(sh * 0.28)

        if content.title:
            tb = slide.shapes.add_textbox(
                title_left, title_top, title_width, int(sh * 0.15),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content.title
            p.alignment = PP_ALIGN.CENTER
            style_runs(p, font_name=style.heading_font, font_size_pt=30,
                        bold=True, color_hex=style.color_text)

        if content.body_paragraphs:
            sub_top = title_top + int(sh * 0.18)
            sub_width = int(sw * 0.6)
            sub_left = (sw - sub_width) // 2
            tb = slide.shapes.add_textbox(
                sub_left, sub_top, sub_width, int(sh * 0.25),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            first = True
            for pd in content.body_paragraphs:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.text = pd.text
                p.alignment = PP_ALIGN.CENTER
                style_runs(p, font_name=style.body_font, font_size_pt=14,
                            color_hex=style.color_muted)

        # Accent line under title
        title_top_pos = int(sh * 0.28)
        line_w = int(sw * 0.08)
        line_left = (sw - line_w) // 2
        line_top = title_top_pos + int(sh * 0.14)
        line_shape = slide.shapes.add_shape(
            1,  # RECTANGLE
            line_left, line_top, line_w, Pt(3),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = rgb(style.color_primary)
        line_shape.line.fill.background()

    # Content images at original positions
    if content.images:
        body_top = int(sh * 0.22)
        _add_content_images(
            slide, style, content.images, body_top,
            has_body_text=bool(content.body_paragraphs),
            preserve_position=bool(content.text_blocks),
        )

    _add_footer(slide, style, slide_number, total_slides, _current_branding)


def _build_section_slide(
    slide, style: TemplateStyle, content: ContentData,
    slide_number: int, total_slides: int,
) -> None:
    """Build a section divider slide."""
    sw, sh = style.slide_width, style.slide_height
    # NOTE: logo already added by build_slide()

    # Large centered section title
    title_width = int(sw * 0.7)
    title_left = (sw - title_width) // 2
    title_top = int(sh * 0.35)

    if content.title:
        tb = slide.shapes.add_textbox(
            title_left, title_top, title_width, int(sh * 0.15),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.alignment = PP_ALIGN.CENTER
        style_runs(p, font_name=style.heading_font, font_size_pt=28,
                    bold=True, color_hex=style.color_text)

    _add_footer(slide, style, slide_number, total_slides, _current_branding)


# ---------------------------------------------------------------------------
# Type-specific renderers
# ---------------------------------------------------------------------------


def _build_agenda_slide(
    slide, style: TemplateStyle, content: ContentData,
    slide_number: int, total_slides: int,
) -> None:
    """Build an agenda/TOC slide with numbered items."""
    sw, sh = style.slide_width, style.slide_height
    margin_left = int(sw * 0.054)

    _add_header(slide, style, "AGENDA")
    _add_footer(slide, style, slide_number, total_slides, _current_branding)

    # Title
    if content.title:
        _add_title_text(slide, style, content.title,
                        margin_left, int(sh * 0.12), int(sw * 0.85))

    # Numbered agenda items — use bold paragraphs as section headers
    body_top = int(sh * 0.22)
    if content.body_paragraphs:
        # Separate into heading items and detail items
        items: list[tuple[str, list[str]]] = []
        current_heading = ""
        current_details: list[str] = []
        for p in content.body_paragraphs:
            if p.bold or (p.font_size >= 14 and len(p.text.split()) <= 8):
                if current_heading:
                    items.append((current_heading, current_details))
                current_heading = p.text
                current_details = []
            elif current_heading:
                current_details.append(p.text)
            else:
                current_heading = p.text
                current_details = []
        if current_heading:
            items.append((current_heading, current_details))

        if items:
            tb = slide.shapes.add_textbox(
                margin_left, body_top, int(sw * 0.88), int(sh * 0.65),
            )
            tf = tb.text_frame
            tf.word_wrap = True

            first = True
            for idx, (heading, details) in enumerate(items):
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                    p.space_before = Pt(8)

                # Numbered heading
                p.text = f"{idx + 1:02d}  {heading}"
                style_runs(p, font_name=style.heading_font, font_size_pt=13,
                            bold=True, color_hex=style.color_text)

                # Sub-details (compact)
                for detail in details[:2]:  # max 2 detail lines per item
                    dp = tf.add_paragraph()
                    dp.text = detail
                    dp.level = 1
                    style_runs(dp, font_name=style.body_font, font_size_pt=10,
                                color_hex=style.color_muted)
        else:
            _add_body_text(slide, style, content.body_paragraphs,
                           margin_left, body_top, int(sw * 0.88), int(sh * 0.65))


def _build_incident_slide(
    slide, style: TemplateStyle, content: ContentData,
    slide_number: int, total_slides: int,
) -> None:
    """Build an incident detail / case study slide with structured layout.

    Layout: title at top, left metadata card, right summary card,
    lower analysis/actions area.
    """
    sw, sh = style.slide_width, style.slide_height
    margin_left = int(sw * 0.054)
    label = _generate_section_label(content)
    _add_header(slide, style, label)
    _add_footer(slide, style, slide_number, total_slides, _current_branding)

    # Title
    if content.title:
        _add_title_text(slide, style, content.title,
                        margin_left, int(sh * 0.12), int(sw * 0.85))

    # Split paragraphs into metadata, summary, and analysis sections
    metadata: list[ParagraphData] = []
    summary: list[ParagraphData] = []
    analysis: list[ParagraphData] = []

    section = "metadata"
    for p in content.body_paragraphs:
        text_lower = p.text.lower()
        if "summary" in text_lower and p.bold:
            section = "summary"
            continue
        if ("analysis" in text_lower or "action" in text_lower) and p.bold:
            section = "analysis"
            continue

        if section == "metadata":
            # Metadata lines contain colons (key: value)
            if ":" in p.text and len(p.text.split()) <= 12:
                metadata.append(p)
            elif p.bold and len(p.text.split()) <= 5:
                section = "summary"  # section header without keyword
                continue
            else:
                metadata.append(p)
        elif section == "summary":
            summary.append(p)
        else:
            analysis.append(p)

    body_top = int(sh * 0.22)
    col_gap = int(sw * 0.02)
    col_w = int(sw * 0.43)

    # Left column: metadata
    if metadata:
        _add_card(slide, style, "Incident Details", metadata,
                  margin_left, body_top, col_w, int(sh * 0.35))

    # Right column: summary + analysis
    right_left = margin_left + col_w + col_gap
    right_paras = summary + analysis
    if right_paras:
        label_text = "Summary & Actions" if analysis else "Summary"
        _add_card(slide, style, label_text, right_paras,
                  right_left, body_top, col_w, int(sh * 0.35))

    # If we have overflow content, add it below
    remaining = [p for p in content.body_paragraphs
                 if p not in metadata and p not in summary and p not in analysis]
    if remaining:
        lower_top = body_top + int(sh * 0.38)
        _add_body_text(slide, style, remaining,
                       margin_left, lower_top, int(sw * 0.88), int(sh * 0.25))

    # Images
    if content.images:
        img_top = body_top + int(sh * 0.38)
        _add_content_images(slide, style, content.images, img_top,
                            has_body_text=True)


def _add_card(
    slide, style: TemplateStyle, card_title: str,
    paragraphs: list[ParagraphData],
    left: int, top: int, width: int, max_height: int,
) -> None:
    """Add a styled card container with title and content."""
    # Card background
    try:
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, max_height,
    )
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(style.color_card)
    card.line.color.rgb = rgb(style.color_line)
    card.line.width = Pt(0.5)

    # Card title
    title_tb = slide.shapes.add_textbox(
        left + Pt(10), top + Pt(6), width - Pt(20), Pt(16),
    )
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.text = card_title
    style_runs(p, font_name=style.heading_font, font_size_pt=10,
                bold=True, color_hex=style.color_primary)

    # Card body
    body_tb = slide.shapes.add_textbox(
        left + Pt(10), top + Pt(24), width - Pt(20), max_height - Pt(30),
    )
    tf = body_tb.text_frame
    tf.word_wrap = True

    # Density-aware sizing
    total_lines = sum(1 + len(p.text) // 50 for p in paragraphs)
    font_pt = 9.0 if total_lines <= 12 else (8.0 if total_lines <= 20 else 7.0)

    first = True
    for pd in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = pd.text
        if pd.bold:
            style_runs(p, font_name=style.heading_font, font_size_pt=font_pt,
                        bold=True, color_hex=style.color_text)
        else:
            style_runs(p, font_name=style.body_font, font_size_pt=font_pt,
                        color_hex=style.color_text)


def _build_kpi_slide(
    slide, style: TemplateStyle, content: ContentData,
    slide_number: int, total_slides: int,
) -> None:
    """Build a KPI / metrics dashboard slide with card grid."""
    sw, sh = style.slide_width, style.slide_height
    margin_left = int(sw * 0.054)
    label = _generate_section_label(content)
    _add_header(slide, style, label)
    _add_footer(slide, style, slide_number, total_slides, _current_branding)

    if content.title:
        _add_title_text(slide, style, content.title,
                        margin_left, int(sh * 0.12), int(sw * 0.85))

    # Group body paragraphs into metric groups (bold heading + detail lines)
    groups: list[tuple[str, list[ParagraphData]]] = []
    current_heading = ""
    current_items: list[ParagraphData] = []

    for p in content.body_paragraphs:
        if p.bold or (p.font_size >= 14 and len(p.text.split()) <= 8):
            if current_heading or current_items:
                groups.append((current_heading, current_items))
            current_heading = p.text
            current_items = []
        else:
            current_items.append(p)
    if current_heading or current_items:
        groups.append((current_heading, current_items))

    if len(groups) >= 2:
        # Card grid layout
        body_top = int(sh * 0.24)
        cols = 2 if len(groups) <= 4 else 3
        col_w = (int(sw * 0.88) - Pt(8) * (cols - 1)) // cols
        rows_needed = (len(groups) + cols - 1) // cols
        row_h = min(int(sh * 0.30), (int(sh * 0.62) - Pt(8) * (rows_needed - 1)) // rows_needed)

        for idx, (heading, items) in enumerate(groups):
            r, c = divmod(idx, cols)
            card_left = margin_left + c * (col_w + Pt(8))
            card_top = body_top + r * (row_h + Pt(8))
            _add_card(slide, style, heading, items,
                      card_left, card_top, col_w, row_h)
    else:
        # Fallback to standard body text
        _add_body_text(slide, style, content.body_paragraphs,
                       margin_left, int(sh * 0.22), int(sw * 0.88), int(sh * 0.65))

    if content.images:
        _add_content_images(slide, style, content.images, int(sh * 0.60),
                            has_body_text=True)


def _build_roadmap_slide(
    slide, style: TemplateStyle, content: ContentData,
    slide_number: int, total_slides: int,
) -> None:
    """Build a roadmap / strategy slide with numbered steps."""
    sw, sh = style.slide_width, style.slide_height
    margin_left = int(sw * 0.054)
    label = _generate_section_label(content)
    _add_header(slide, style, label)
    _add_footer(slide, style, slide_number, total_slides, _current_branding)

    if content.title:
        _add_title_text(slide, style, content.title,
                        margin_left, int(sh * 0.12), int(sw * 0.85))

    # Split into intro paragraph + action items
    intro: list[ParagraphData] = []
    steps: list[tuple[str, list[ParagraphData]]] = []
    current_step = ""
    current_detail: list[ParagraphData] = []

    for p in content.body_paragraphs:
        if p.bold and len(p.text.split()) <= 10:
            if current_step:
                steps.append((current_step, current_detail))
            elif not steps and current_detail:
                intro = current_detail[:]
            current_step = p.text
            current_detail = []
        else:
            current_detail.append(p)
    if current_step:
        steps.append((current_step, current_detail))
    elif current_detail and not intro:
        intro = current_detail

    body_top = int(sh * 0.22)

    # Intro paragraph
    if intro:
        _add_body_text(slide, style, intro,
                       margin_left, body_top, int(sw * 0.88), int(sh * 0.10))
        body_top += int(sh * 0.12)

    # Steps as cards
    if len(steps) >= 2:
        cols = 2 if len(steps) <= 4 else 3
        col_w = (int(sw * 0.88) - Pt(8) * (cols - 1)) // cols
        rows_needed = (len(steps) + cols - 1) // cols
        avail_h = int(sh * 0.87) - body_top
        row_h = min(int(sh * 0.25), (avail_h - Pt(8) * (rows_needed - 1)) // rows_needed)

        for idx, (heading, items) in enumerate(steps):
            r, c = divmod(idx, cols)
            card_left = margin_left + c * (col_w + Pt(8))
            card_top = body_top + r * (row_h + Pt(8))
            _add_card(slide, style, heading, items,
                      card_left, card_top, col_w, row_h)
    elif content.body_paragraphs:
        _add_body_text(slide, style, content.body_paragraphs,
                       margin_left, body_top, int(sw * 0.88), int(sh * 0.65))

    if content.images:
        _add_content_images(slide, style, content.images, body_top,
                            has_body_text=True)


def _build_generic_content_slide(
    slide, style: TemplateStyle, content: ContentData,
    slide_number: int, total_slides: int,
) -> None:
    """Build a standard content slide (narrative, bullets, etc.)."""
    sw, sh = style.slide_width, style.slide_height
    margin_left = int(sw * 0.054)
    label = _generate_section_label(content)
    _add_header(slide, style, label)
    _add_footer(slide, style, slide_number, total_slides, _current_branding)

    title_top = int(sh * 0.12)
    body_top = int(sh * 0.22)
    body_max_h = int(sh * 0.65)

    if content.title:
        _add_title_text(slide, style, content.title,
                        margin_left, title_top, int(sw * 0.85))

    has_images = bool(content.images)

    if content.text_blocks:
        _add_text_blocks(slide, style, content.text_blocks)
    elif content.body_paragraphs:
        bw = int(sw * 0.50) if has_images else int(sw * 0.88)
        _add_body_text(slide, style, content.body_paragraphs,
                       margin_left, body_top, bw, body_max_h)

    # Tables
    if content.tables:
        table_top = body_top
        if content.body_paragraphs and not content.text_blocks:
            est_lines = sum(1 + len(p.text) // 60 for p in content.body_paragraphs)
            table_top = body_top + int(est_lines * Pt(16))
            table_top = min(table_top, int(sh * 0.55))
        for td in content.tables:
            _add_table(slide, style, td["data"],
                       margin_left, table_top, int(sw * 0.85), int(sh * 0.40))
            break

    if content.images:
        _add_content_images(slide, style, content.images, body_top,
                            has_body_text=bool(content.body_paragraphs),
                            preserve_position=bool(content.text_blocks))


# ---------------------------------------------------------------------------
# Slide type → renderer dispatch
# ---------------------------------------------------------------------------

_TYPE_RENDERERS: dict[str, callable] = {}


def _get_renderer(slide_type: str):
    """Return the appropriate renderer for a given slide type."""
    dispatch = {
        "title": _build_title_slide,
        "section": _build_section_slide,
        "agenda": _build_agenda_slide,
        "toc": _build_agenda_slide,
        "closing": _build_generic_content_slide,
        "metrics_dashboard": _build_kpi_slide,
        "comparison": _build_kpi_slide,
        "process_flow": _build_roadmap_slide,
        "timeline": _build_roadmap_slide,
        "content_narrative": _build_generic_content_slide,
        "content_bullets": _build_generic_content_slide,
        "image_heavy": _build_generic_content_slide,
        "data": _build_kpi_slide,
    }
    return dispatch.get(slide_type, _build_generic_content_slide)


def _is_incident_slide(content: ContentData) -> bool:
    """Heuristic: detect incident/case-study slides by content pattern."""
    title_lower = (content.title or "").lower()
    if "incident" in title_lower or "case study" in title_lower:
        return True
    # Check for metadata-heavy content (lots of "Key: Value" lines)
    # Require a high count to avoid false positives on content slides
    # that happen to have a few labeled fields
    kv_count = sum(1 for p in content.body_paragraphs
                   if ":" in p.text and len(p.text.split()) <= 12)
    return kv_count >= 8


# ---------------------------------------------------------------------------
# Main slide builder (public)
# ---------------------------------------------------------------------------


def build_slide(
    prs: Presentation, style: TemplateStyle,
    content: ContentData, slide_number: int, total_slides: int,
    branding: BrandingPolicy | None = None,
) -> None:
    """Build a single output slide from scratch."""
    global _current_branding
    _current_branding = branding or BrandingPolicy()
    blank_layout = _find_blank_layout(prs)
    slide = prs.slides.add_slide(blank_layout)

    # Background
    _add_background(slide, style)

    # Decorative shapes (skip for title/section slides for cleaner look)
    if content.slide_type not in ("title", "section"):
        _add_decorative_shapes(slide, style)

    # Logo
    _add_logo(slide, style)

    # Dispatch to type-specific renderer
    if _is_incident_slide(content):
        renderer = _build_incident_slide
    else:
        renderer = _get_renderer(content.slide_type)

    renderer(slide, style, content, slide_number, total_slides)

    # Speaker notes
    if content.notes:
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            if tf:
                tf.text = content.notes
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Orchestrator (public)
# ---------------------------------------------------------------------------


def apply_recreate(
    template_path: Path, content_path: Path, output_path: Path,
    config: TransferConfig,
) -> dict[str, Any]:
    """Recreate mode: analyze template style, extract content, rebuild from scratch."""
    from pptx_template_transfer.analysis.theme_extractor import analyze_template
    from pptx_template_transfer.extraction.content_extractor import extract_all_content
    from pptx_template_transfer.transform.layout_mapper import map_content_to_layout

    report: dict[str, Any] = {
        "mode": "recreate",
        "slides": [],
        "warnings": [],
        "errors": [],
    }
    th = config.thresholds

    # Step 1: Analyze template
    print("\n[recreate] Analyzing template style...")
    style = analyze_template(template_path)
    print(f"  Fonts: heading={style.heading_font}, body={style.body_font}")
    print(
        f"  Colors: primary=#{style.color_primary}, "
        f"text=#{style.color_text}, bg=#{style.color_background}"
    )
    print(
        f"  Logo: {'found' if style.logo_blob else 'not found'} "
        f"({style.logo_width}x{style.logo_height})"
    )
    print(f"  Footer: '{style.footer_company}'")

    # Step 2: Extract content
    print("\n[recreate] Extracting content...")
    content_list = extract_all_content(content_path, th)
    ct = len(content_list)

    # Check aspect ratio compatibility — if content and template differ
    # significantly (e.g. 4:3 vs 16:9), text_blocks positions won't transfer
    # correctly, so fall back to reflowed body_paragraphs layout.
    content_prs = Presentation(str(content_path))
    content_ratio = (
        content_prs.slide_width / content_prs.slide_height
        if content_prs.slide_height
        else 1.0
    )
    template_ratio = (
        style.slide_width / style.slide_height
        if style.slide_height
        else 1.0
    )
    ratio_diff = abs(content_ratio - template_ratio) / max(
        content_ratio, template_ratio
    )
    use_text_blocks = ratio_diff < 0.15  # Within 15% aspect ratio = safe

    if not use_text_blocks:
        print(
            f"  Aspect ratio mismatch: content={content_ratio:.2f}, "
            f"template={template_ratio:.2f} (diff={ratio_diff:.0%}) "
            f"— using reflowed layout instead of position-preserving"
        )
        for cd in content_list:
            cd.text_blocks = []  # Force flat layout
    else:
        print(
            f"  Aspect ratios compatible ({ratio_diff:.0%} diff) "
            f"— preserving positions"
        )

    # Step 2b: Layout-aware content placement
    # Use the layout mapper to get zone assignments for smarter placement.
    try:
        layout_assignments = map_content_to_layout(content_list, style)
        if layout_assignments and config.verbose:
            print(f"  Layout mapper assigned zones for {len(layout_assignments)} slides")
    except Exception as exc:
        layout_assignments = None
        log.debug("Layout mapper unavailable or failed: %s", exc)

    for i, cd in enumerate(content_list):
        if config.verbose:
            print(
                f"  Slide {i+1}: type={cd.slide_type}, words={cd.word_count}, "
                f"title='{cd.title[:40]}', paras={len(cd.body_paragraphs)}, "
                f"tables={len(cd.tables)}, images={len(cd.images)}"
            )

    # Step 3: Build output — use template as base to preserve theme/masters
    print(f"\n[recreate] Building {ct} slides from scratch...")
    output_prs = Presentation(str(template_path))

    # Remove all template slides but keep masters/theme
    sld_id_lst = output_prs.slides._sldIdLst
    ns_r = NSMAP["r"]
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(f"{{{ns_r}}}id")
        if r_id:
            try:
                output_prs.part.drop_rel(r_id)
            except Exception:
                pass
        sld_id_lst.remove(sld_id)

    for i, cd in enumerate(content_list):
        slide_report: dict[str, Any] = {
            "index": i + 1,
            "source_slide": cd.source_slide_index + 1,
            "content_type": cd.slide_type,
            "title": cd.title[:80] if cd.title else "",
            "word_count": cd.word_count,
            "status": "ok",
            "provenance": {
                "title": "source_content",
                "body": "source_content",
                "footer": "target_shell",
                "section_label": "converter_generated_bridge",
            },
        }
        try:
            build_slide(output_prs, style, cd, i + 1, ct, config.branding)
            print(
                f"  Slide {i+1}/{ct}: [{cd.slide_type}] "
                f'"{cd.title[:50] if cd.title else "(no title)"}"'
            )
        except Exception as exc:
            slide_report["status"] = "error"
            slide_report["error"] = str(exc)
            report["errors"].append(f"Slide {i+1}: {exc}")
            log.error(
                "Slide %d failed: %s\n%s", i + 1, exc, traceback.format_exc()
            )
            print(f"  Slide {i+1}/{ct}: ERROR - {exc}")
            # Add blank slide as fallback
            try:
                output_prs.slides.add_slide(_find_blank_layout(output_prs))
            except Exception:
                pass
        report["slides"].append(slide_report)

    # Save
    print(f"\n[recreate] Saving to {output_path}...")
    output_prs.save(str(output_path))
    success = sum(1 for s in report["slides"] if s["status"] == "ok")
    print(f"[recreate] Done! {success}/{ct} slides created successfully.")

    # --- Post-generation validation ---
    from pptx_template_transfer.validation.contamination_checker import (
        check_target_contamination,
    )
    from pptx_template_transfer.validation.source_coverage import (
        compute_source_coverage,
    )

    output_prs_check = Presentation(str(output_path))

    # Target contamination check
    contamination = check_target_contamination(output_prs_check, template_path)
    if contamination:
        report["warnings"].extend(contamination)
        for w in contamination:
            print(f"  WARNING: {w}")

    # Source coverage check
    coverage = compute_source_coverage(output_prs_check, content_list)
    report["source_coverage"] = {
        "overall_pct": coverage.overall_text_coverage_pct,
        "total_source_slides": coverage.total_source_slides,
        "total_output_slides": coverage.total_output_slides,
        "unmapped": coverage.unmapped_source_slides,
        "entries": [
            {
                "source_slide": e.source_slide_index + 1,
                "title": e.source_title,
                "text_used_pct": e.text_used_pct,
                "output_slides": [x + 1 for x in e.output_slide_indices],
                "tables_dropped": e.tables_dropped,
                "images_dropped": e.images_dropped,
                "charts_dropped": e.charts_dropped,
            }
            for e in coverage.entries
        ],
    }
    if coverage.warnings:
        report["warnings"].extend(coverage.warnings)
        for w in coverage.warnings:
            print(f"  COVERAGE: {w}")

    print(f"  Source coverage: {coverage.overall_text_coverage_pct}%")

    # Quality report
    from pptx_template_transfer.validation.quality_report import (
        generate_quality_report,
    )
    qr = generate_quality_report(output_prs_check, content_list, config)
    report["quality"] = {
        "overall_score": round(qr.overall_score, 1),
        "native_count": qr.native_count,
        "fallback_count": qr.fallback_count,
        "slides_needing_review": [
            {
                "slide": sq.slide_index + 1,
                "score": round(getattr(sq, "slide_score", sq.content_coverage_pct), 1),
                "reasons": sq.review_reasons,
            }
            for sq in qr.slides if sq.needs_manual_review
        ],
    }
    if qr.warnings:
        report["warnings"].extend(qr.warnings)
    print(f"  Quality score: {qr.overall_score:.1f}/100")
    if qr.warnings:
        for w in qr.warnings[:5]:
            print(f"  QUALITY: {w}")

    if config.report_path:
        clean = json.loads(json.dumps(report, default=str))
        config.report_path.write_text(json.dumps(clean, indent=2))
        print(f"Report written to {config.report_path}")

    return report
