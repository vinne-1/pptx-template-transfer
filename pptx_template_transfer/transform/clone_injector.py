"""Clone/design mode pipeline: slide matching, cloning, text injection, and orchestration."""
from __future__ import annotations

import io
import logging
import math
import traceback
from copy import deepcopy
from datetime import date
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu, Pt

from pptx_template_transfer.analysis.slide_classifier import (
    classify_all_shapes,
    classify_template_structure,
)
from pptx_template_transfer.extraction.content_extractor import extract_content
from pptx_template_transfer.helpers import (
    DATE_PATTERN,
    FOOTER_PATTERNS,
    JUST_NUMBER_RE,
    NSMAP,
    PAGE_NUM_PATTERN,
    PH_DATE,
    PH_FOOTER_SET,
    PH_SLIDE_NUMBER,
    has_placeholder_type,
    is_chart,
    is_group,
    is_ole_or_embedded,
    is_picture,
    is_table,
    max_font_pt,
    placeholder_type_int,
    shape_area,
    shape_area_pct,
    shape_bottom_frac,
    shape_top_frac,
    text_of,
    update_rids_in_tree,
    word_count,
)
from pptx_template_transfer.models import (
    ContentData,
    ParagraphData,
    Thresholds,
    TransferConfig,
)

log = logging.getLogger("pptx_template_transfer")


# ============================================================================
# SMART SLIDE MATCHING
# ============================================================================

_TYPE_COMPAT = {
    ("title", "title"): 40, ("title", "section"): 20, ("title", "narrative"): 10,
    ("content", "narrative"): 40, ("content", "list"): 30, ("content", "grid"): 25,
    ("content", "data"): 15, ("content", "title"): 8,
    ("section", "section"): 40, ("section", "title"): 25,
    ("data", "data"): 40, ("data", "narrative"): 20, ("data", "grid"): 25,
    ("closing", "closing"): 40, ("closing", "section"): 20,
    ("image", "visual"): 40, ("image", "narrative"): 15,
    ("blank", "section"): 10,
}


def _match_score(
    c_type: str, t_struct: str,
    ci: int, ti: int, ct: int, tt: int,
    c_words: int, t_words: int,
    c_has_table: bool, t_has_table: bool,
    c_paras: int, t_is_list: bool,
) -> float:
    score = float(_TYPE_COMPAT.get((c_type, t_struct), 5))

    if t_words > 0 and c_words > 0:
        score += 25 * min(c_words, t_words) / max(c_words, t_words)
    elif c_words == 0 and t_words <= 10:
        score += 20

    if c_has_table and t_has_table:
        score += 20
    elif c_has_table:
        score += 5
    elif c_paras >= 5 and t_is_list:
        score += 15
    elif c_paras >= 3 and t_struct == "narrative":
        score += 15
    elif c_paras < 3 and t_struct in ("section", "title"):
        score += 10

    if ci == 0 and t_struct == "title":
        score += 15
    elif ci == ct - 1 and t_struct == "closing":
        score += 15
    elif ct > 1 and tt > 1:
        score += 15 * (1 - abs(ci / (ct - 1) - ti / (tt - 1)))

    return score


def build_slide_mapping(
    content_prs: Presentation, template_prs: Presentation,
    content_data_list: list[ContentData], th: Thresholds,
) -> list[int]:
    sw, sh = template_prs.slide_width, template_prs.slide_height
    ct, tt = len(content_prs.slides), len(template_prs.slides)

    t_info = []
    for i, slide in enumerate(template_prs.slides):
        struct = classify_template_structure(slide, sw, sh, i, tt)
        words = sum(word_count(text_of(s)) for s in slide.shapes)
        t_info.append({
            "struct": struct, "words": words,
            "has_table": any(is_table(s) for s in slide.shapes),
            "is_list": struct in ("list", "grid"),
        })

    # Score matrix
    score_matrix: list[list[tuple[int, float]]] = []
    for ci, cd in enumerate(content_data_list):
        scores = []
        for ti, tinfo in enumerate(t_info):
            sc = _match_score(
                cd.slide_type, tinfo["struct"], ci, ti, ct, tt,
                cd.word_count, tinfo["words"],
                len(cd.tables) > 0, tinfo["has_table"],
                len(cd.body_paragraphs), tinfo["is_list"],
            )
            scores.append((ti, sc))
        scores.sort(key=lambda x: -x[1])
        score_matrix.append(scores)

    # Greedy with variety
    usage: dict[int, int] = {i: 0 for i in range(tt)}
    max_per = max(2, math.ceil(ct * th.variety_max_pct))
    min_distinct = min(tt, max(3, math.ceil(ct / 3)))

    mapping = []
    for ci, scores in enumerate(score_matrix):
        best_idx, best_sc = scores[0]
        if usage[best_idx] >= max_per:
            for ti2, sc2 in scores[1:]:
                if usage[ti2] < max_per:
                    best_idx = ti2
                    break
        mapping.append(best_idx)
        usage[best_idx] += 1

    # Redistribute to hit min_distinct
    used_set = {ti for ti, c in usage.items() if c > 0}
    unused = [ti for ti in range(tt) if usage[ti] == 0]
    if len(used_set) < min_distinct and unused:
        overused = sorted(
            [(ti, c) for ti, c in usage.items() if c > 1], key=lambda x: -x[1],
        )
        for u_ti in unused:
            if not overused:
                break
            donor_ti = overused[0][0]
            candidates = [(ci, score_matrix[ci]) for ci in range(ct) if mapping[ci] == donor_ti]
            best_ci, best_sc = None, -1.0
            for ci2, scores in candidates:
                for ti2, sc2 in scores:
                    if ti2 == u_ti and sc2 > best_sc:
                        best_ci, best_sc = ci2, sc2
            if best_ci is not None and best_sc > 10:
                usage[mapping[best_ci]] -= 1
                mapping[best_ci] = u_ti
                usage[u_ti] = usage.get(u_ti, 0) + 1
                overused = sorted(
                    [(ti, c) for ti, c in usage.items() if c > 1], key=lambda x: -x[1],
                )

    for ci in range(ct):
        sc_val = next((sc for t, sc in score_matrix[ci] if t == mapping[ci]), 0)
        log.debug(
            "  Slide %d (%s, %dw) -> Template %d (%s) score=%.0f",
            ci + 1, content_data_list[ci].slide_type,
            content_data_list[ci].word_count, mapping[ci] + 1,
            t_info[mapping[ci]]["struct"], sc_val,
        )
    used = sum(1 for v in usage.values() if v > 0)
    log.debug("  Variety: %d/%d templates used (target: %d+)", used, tt, min_distinct)

    return mapping


# ============================================================================
# SLIDE CLONING
# ============================================================================

def _clone_slide(template_prs: Presentation, src_slide, dst_prs: Presentation):
    dst_layout = dst_prs.slide_layouts[0]
    src_layout_name = src_slide.slide_layout.name
    for layout in dst_prs.slide_layouts:
        if layout.name == src_layout_name:
            dst_layout = layout
            break

    new_slide = dst_prs.slides.add_slide(dst_layout)

    spTree = new_slide.shapes._spTree
    for sp in list(spTree):
        tag = etree.QName(sp.tag).localname if isinstance(sp.tag, str) else ""
        if tag in ("sp", "pic", "grpSp", "cxnSp", "graphicFrame"):
            spTree.remove(sp)

    src_spTree = src_slide.shapes._spTree
    for child in src_spTree:
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""
        if tag in ("sp", "pic", "grpSp", "cxnSp", "graphicFrame"):
            spTree.append(deepcopy(child))

    # Background
    src_sld, dst_sld = src_slide._element, new_slide._element
    src_bg = src_sld.find(f'{{{NSMAP["p"]}}}bg')
    if src_bg is not None:
        dst_bg = dst_sld.find(f'{{{NSMAP["p"]}}}bg')
        if dst_bg is not None:
            dst_sld.remove(dst_bg)
        new_bg = deepcopy(src_bg)
        cSld = dst_sld.find(f'{{{NSMAP["p"]}}}cSld')
        if cSld is not None:
            dst_sld.insert(list(dst_sld).index(cSld), new_bg)
        else:
            dst_sld.insert(0, new_bg)

    # Transition
    ns_p = NSMAP["p"]
    src_transition = src_sld.find(f'{{{ns_p}}}transition')
    if src_transition is not None:
        dst_sld.append(deepcopy(src_transition))

    # Relationships
    rid_map: dict[str, str] = {}
    broken_rels: list[str] = []
    for rel_key, rel in src_slide.part.rels.items():
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
            continue
        try:
            if rel.is_external:
                new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rid = new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
            rid_map[rel_key] = new_rid
        except Exception as exc:
            broken_rels.append(f"{rel_key} ({rel.reltype}): {exc}")

    if broken_rels:
        log.warning("  Broken relationships during clone: %s", "; ".join(broken_rels))

    if rid_map:
        update_rids_in_tree(spTree, rid_map)
        dst_bg2 = dst_sld.find(f'{{{NSMAP["p"]}}}bg')
        if dst_bg2 is not None:
            update_rids_in_tree(dst_bg2, rid_map)

    return new_slide


# ============================================================================
# TEXT INJECTION - multi-level format preservation
# ============================================================================

def _save_all_paragraph_formats(shape) -> dict[int, tuple]:
    """Save formatting from all paragraph levels in a shape.

    Returns dict mapping indent level -> (pPr_element, rPr_element).
    """
    ns_a = NSMAP["a"]
    formats: dict[int, tuple] = {}
    if not shape.has_text_frame:
        return formats

    for para in shape.text_frame.paragraphs:
        p_el = para._p
        level = para.level if para.level else 0
        if level in formats:
            continue

        pPr = p_el.find(f'{{{ns_a}}}pPr')
        pPr_copy = deepcopy(pPr) if pPr is not None else None

        rPr_copy = None
        for r in p_el.findall(f'{{{ns_a}}}r'):
            rPr = r.find(f'{{{ns_a}}}rPr')
            if rPr is not None:
                rPr_copy = deepcopy(rPr)
                break
        if rPr_copy is None:
            endRPr = p_el.find(f'{{{ns_a}}}endParaRPr')
            if endRPr is not None:
                rPr_copy = deepcopy(endRPr)

        formats[level] = (pPr_copy, rPr_copy)

    return formats


def _get_format_for_level(formats: dict[int, tuple], level: int) -> tuple:
    """Get (pPr, rPr) for a given indent level, falling back to closest."""
    if level in formats:
        return formats[level]
    if not formats:
        return (None, None)
    closest = min(formats.keys(), key=lambda k: abs(k - level))
    return formats[closest]


def _estimate_text_capacity(shape, slide_w: int, slide_h: int, th: Thresholds) -> int:
    """Estimate how many characters a text shape can hold."""
    w_inches = (shape.width or 0) / 914400.0
    h_inches = (shape.height or 0) / 914400.0
    area_sq_inches = w_inches * h_inches
    if area_sq_inches <= 0:
        return 100
    return max(20, int(area_sq_inches * th.overflow_chars_per_sq_inch))


def _fit_paragraphs(
    paragraphs: list[ParagraphData], max_chars: int,
) -> list[ParagraphData]:
    """Truncate paragraphs to fit within max_chars, adding '...' if needed."""
    result = []
    chars_used = 0
    for p in paragraphs:
        if chars_used + len(p.text) <= max_chars:
            result.append(p)
            chars_used += len(p.text)
        else:
            remaining = max_chars - chars_used
            if remaining > 20:
                truncated = ParagraphData(
                    text=p.text[:remaining - 3] + "...",
                    level=p.level, bold=p.bold, italic=p.italic,
                    font_size=p.font_size, runs=p.runs,
                )
                result.append(truncated)
            elif not result:
                # At least include one truncated paragraph
                result.append(ParagraphData(
                    text=p.text[:max(50, max_chars)] + "...",
                    level=p.level, bold=p.bold,
                ))
            break
    return result


def _inject_text_simple(shape, text: str) -> None:
    """Replace text preserving first paragraph's formatting."""
    if not shape.has_text_frame or not shape.text_frame.paragraphs:
        return
    ns_a = NSMAP["a"]
    formats = _save_all_paragraph_formats(shape)
    pPr, rPr = _get_format_for_level(formats, 0)

    txBody = shape.text_frame._txBody
    for p in list(txBody.findall(f'{{{ns_a}}}p')):
        txBody.remove(p)

    for para_text in text.split("\n"):
        new_p = etree.SubElement(txBody, f'{{{ns_a}}}p')
        if pPr is not None:
            new_p.append(deepcopy(pPr))
        if para_text.strip():
            new_r = etree.SubElement(new_p, f'{{{ns_a}}}r')
            if rPr is not None:
                new_r.append(deepcopy(rPr))
            new_t = etree.SubElement(new_r, f'{{{ns_a}}}t')
            new_t.text = para_text
        else:
            eRPr = etree.SubElement(new_p, f'{{{ns_a}}}endParaRPr')
            if rPr is not None:
                for k, v in rPr.attrib.items():
                    eRPr.attrib[k] = v


def _inject_structured_text(
    shape, paragraphs: list[ParagraphData], th: Thresholds,
    slide_w: int = 0, slide_h: int = 0,
) -> None:
    """Inject structured paragraphs with multi-level format preservation."""
    if not shape.has_text_frame or not paragraphs:
        return
    if not shape.text_frame.paragraphs:
        return

    ns_a = NSMAP["a"]

    # Save formatting per level
    formats = _save_all_paragraph_formats(shape)

    # Overflow prevention
    if slide_w > 0 and slide_h > 0:
        capacity = _estimate_text_capacity(shape, slide_w, slide_h, th)
        total_chars = sum(len(p.text) for p in paragraphs)
        if total_chars > capacity:
            paragraphs = _fit_paragraphs(paragraphs, capacity)
            log.debug("    Overflow: truncated %d -> %d chars (capacity=%d)",
                       total_chars, sum(len(p.text) for p in paragraphs), capacity)

    txBody = shape.text_frame._txBody
    for p in list(txBody.findall(f'{{{ns_a}}}p')):
        txBody.remove(p)

    for pd in paragraphs:
        new_p = etree.SubElement(txBody, f'{{{ns_a}}}p')
        pPr_tmpl, rPr_tmpl = _get_format_for_level(formats, pd.level)

        if pPr_tmpl is not None:
            pPr = deepcopy(pPr_tmpl)
            if pd.level > 0:
                pPr.set("lvl", str(pd.level))
            new_p.append(pPr)
        elif pd.level > 0:
            pPr = etree.SubElement(new_p, f'{{{ns_a}}}pPr')
            pPr.set("lvl", str(pd.level))

        if pd.text.strip():
            # If we have runs with hyperlinks, use multi-run injection
            if pd.runs and any(r.hyperlink_url for r in pd.runs):
                for rd in pd.runs:
                    if not rd.text:
                        continue
                    new_r = etree.SubElement(new_p, f'{{{ns_a}}}r')
                    if rPr_tmpl is not None:
                        rPr = deepcopy(rPr_tmpl)
                        if pd.bold or rd.bold:
                            rPr.set("b", "1")
                        if pd.italic or rd.italic:
                            rPr.set("i", "1")
                        # Hyperlink
                        if rd.hyperlink_url:
                            hlinkClick = etree.SubElement(
                                rPr, f'{{{ns_a}}}hlinkClick',
                            )
                            hlinkClick.set(
                                f'{{{NSMAP["r"]}}}id', "",  # Will need rel
                            )
                        new_r.append(rPr)
                    new_t = etree.SubElement(new_r, f'{{{ns_a}}}t')
                    new_t.text = rd.text
            else:
                new_r = etree.SubElement(new_p, f'{{{ns_a}}}r')
                if rPr_tmpl is not None:
                    rPr = deepcopy(rPr_tmpl)
                    if pd.bold:
                        rPr.set("b", "1")
                    if pd.italic:
                        rPr.set("i", "1")
                    new_r.append(rPr)
                new_t = etree.SubElement(new_r, f'{{{ns_a}}}t')
                new_t.text = pd.text
        else:
            eRPr = etree.SubElement(new_p, f'{{{ns_a}}}endParaRPr')
            if rPr_tmpl is not None:
                for k, v in rPr_tmpl.attrib.items():
                    eRPr.attrib[k] = v


def _clear_shape_text(shape) -> None:
    if not shape.has_text_frame:
        return
    ns_a = NSMAP["a"]
    txBody = shape.text_frame._txBody
    for p in list(txBody.findall(f'{{{ns_a}}}p')):
        txBody.remove(p)
    new_p = etree.SubElement(txBody, f'{{{ns_a}}}p')
    etree.SubElement(new_p, f'{{{ns_a}}}endParaRPr')


# ---- Template text clearing: is_protected + prepare_cloned_slide ----

def _is_protected_shape(shape, slide_w: int, slide_h: int) -> bool:
    """Return True if the shape should keep its template text untouched.

    Aggressively clears ALL template text except truly structural elements
    (footers, page numbers, dates, confidential notices, empty shapes).
    This ensures zero template words leak into the output.
    """
    # No text frame - nothing to clear
    if not shape.has_text_frame:
        return True

    # Media shapes (picture, chart, table, group, OLE)
    if is_picture(shape) or is_chart(shape) or is_table(shape) or is_group(shape):
        return True
    if is_ole_or_embedded(shape):
        return True

    # Empty text - nothing to clear
    text = shape.text_frame.text.strip()
    if not text:
        return True

    # Placeholder-based footer (slide number, date, footer) - structural
    ph = placeholder_type_int(shape)
    if ph is not None and ph in PH_FOOTER_SET:
        return True

    # Footer zone (bottom 8% of slide) - structural
    bottom_frac = shape_bottom_frac(shape, slide_h)
    if bottom_frac > 0.92:
        return True

    # Common footer/label patterns (Page XX, Confidential, dates, (c))
    if FOOTER_PATTERNS.match(text.strip()):
        return True

    # Just a number like "01", "02"
    if JUST_NUMBER_RE.match(text.strip()):
        return True

    # NOT protected - this is an injection target whose text gets erased
    return False


def _prepare_cloned_slide(
    slide, slide_w: int, slide_h: int,
) -> tuple[list, list]:
    """Erase template text from injection targets, leave protected shapes untouched.

    Returns (injection_targets, protected_shapes) for diagnostic tracking.
    """
    targets: list = []
    protected: list = []

    for shape in slide.shapes:
        if _is_protected_shape(shape, slide_w, slide_h):
            protected.append(shape)
            continue

        # This is an injection target - CLEAR its text
        targets.append(shape)
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = ""

    return targets, protected


# ============================================================================
# CONTENT INJECTION
# ============================================================================

def inject_content(
    cloned_slide, content_data: ContentData,
    slide_w: int, slide_h: int, th: Thresholds,
) -> dict[str, Any]:
    """Inject content into a cloned slide. Returns diagnostic dict."""
    diag: dict[str, Any] = {
        "shapes": [], "injected_title": None,
        "injected_body": None, "protected_count": 0,
        "cleared_count": 0,
    }

    # Step 1: Classify shapes BEFORE clearing (needs original text for accuracy)
    classifications = classify_all_shapes(cloned_slide, slide_w, slide_h, th)

    # Step 2: Determine which shapes are injection targets vs protected
    targets, protected = _prepare_cloned_slide(cloned_slide, slide_w, slide_h)
    target_ids = {id(s) for s in targets}
    diag["cleared_count"] = len(targets)
    diag["protected_count"] = len(protected)

    # Step 3: Use pre-clearing classifications to assign title/body zones
    title_shape = None
    body_shapes: list = []

    for shape, role, conf in classifications:
        diag["shapes"].append({
            "name": shape.name, "role": role,
            "confidence": round(conf, 2),
            "area_pct": round(shape_area_pct(shape, slide_w, slide_h), 1),
            "top_pct": round(shape_top_frac(shape, slide_h) * 100, 0),
            "text_preview": text_of(shape)[:40],
            "is_target": id(shape) in target_ids,
        })
        # Only assign zones from shapes that are injection targets
        if id(shape) not in target_ids:
            continue
        if role == "title" and title_shape is None:
            title_shape = shape
        elif role in ("body", "info"):
            body_shapes.append(shape)

    # Fallback: if classifier didn't find title/body among targets,
    # pick from targets by font size (title) and area (body)
    if not title_shape and targets:
        top_half = [s for s in targets
                    if s.has_text_frame and shape_top_frac(s, slide_h) < 0.45]
        if top_half:
            title_shape = max(top_half, key=lambda s: max_font_pt(s))

    if not body_shapes:
        for s in sorted(targets, key=lambda s: shape_area(s), reverse=True):
            if s != title_shape and s.has_text_frame:
                body_shapes.append(s)
            if len(body_shapes) >= th.body_max_zones:
                break

    # Step 4: Inject content into zones (shapes already cleared by step 2)
    # --- Title ---
    if content_data.title and title_shape:
        _inject_text_simple(title_shape, content_data.title)
        diag["injected_title"] = content_data.title[:50]

    # --- Body ---
    if content_data.body_paragraphs and body_shapes:
        if len(body_shapes) == 1:
            _inject_structured_text(
                body_shapes[0], content_data.body_paragraphs, th, slide_w, slide_h,
            )
            wc = sum(word_count(p.text) for p in content_data.body_paragraphs)
            diag["injected_body"] = f"{wc} words -> 1 zone"
        else:
            per_zone = max(1, len(content_data.body_paragraphs) // len(body_shapes))
            idx = 0
            for i, zone in enumerate(body_shapes):
                chunk = (content_data.body_paragraphs[idx:]
                         if i == len(body_shapes) - 1
                         else content_data.body_paragraphs[idx:idx + per_zone])
                idx += per_zone
                if chunk:
                    _inject_structured_text(zone, chunk, th, slide_w, slide_h)
            wc = sum(word_count(p.text) for p in content_data.body_paragraphs)
            diag["injected_body"] = f"{wc} words -> {len(body_shapes)} zones"

    return diag


# ============================================================================
# TABLE, CHART & IMAGE HANDLING
# ============================================================================

def _inject_table_cell_text(cell, text: str) -> None:
    """Fill a table cell preserving its formatting."""
    ns_a = NSMAP["a"]
    tf = cell.text_frame
    if not tf.paragraphs:
        cell.text = text
        return
    # Save format from first paragraph
    first_p = tf.paragraphs[0]._p
    rPr = None
    for r in first_p.findall(f'{{{ns_a}}}r'):
        rp = r.find(f'{{{ns_a}}}rPr')
        if rp is not None:
            rPr = deepcopy(rp)
            break
    pPr = first_p.find(f'{{{ns_a}}}pPr')
    pPr = deepcopy(pPr) if pPr is not None else None

    # Clear and refill
    txBody = tf._txBody
    for p in list(txBody.findall(f'{{{ns_a}}}p')):
        txBody.remove(p)
    new_p = etree.SubElement(txBody, f'{{{ns_a}}}p')
    if pPr is not None:
        new_p.append(pPr)
    new_r = etree.SubElement(new_p, f'{{{ns_a}}}r')
    if rPr is not None:
        new_r.append(rPr)
    new_t = etree.SubElement(new_r, f'{{{ns_a}}}t')
    new_t.text = text


def _add_table_rows(table, count: int) -> None:
    """Clone the last row of a table to add more rows."""
    ns_a = NSMAP["a"]
    tbl_el = table._tbl
    rows = tbl_el.findall(f'{{{ns_a}}}tr')
    if not rows:
        return
    last_row = rows[-1]
    for _ in range(count):
        new_row = deepcopy(last_row)
        # Clear text in cloned cells
        for tc in new_row.findall(f'{{{ns_a}}}tc'):
            for p in tc.findall(f'.//{{{ns_a}}}p'):
                for r in p.findall(f'{{{ns_a}}}r'):
                    t = r.find(f'{{{ns_a}}}t')
                    if t is not None:
                        t.text = ""
        tbl_el.append(new_row)


def _handle_tables(cloned_slide, content_data: ContentData, slide_w: int, slide_h: int) -> None:
    if not content_data.tables:
        return

    template_tables = [s for s in cloned_slide.shapes if is_table(s)]
    ct = content_data.tables[0]

    if template_tables:
        tmpl = template_tables[0].table
        c_data = ct["data"]
        if not c_data:
            return

        c_rows, c_cols = len(c_data), len(c_data[0]) if c_data else 0
        t_rows, t_cols = len(tmpl.rows), len(tmpl.columns)

        # Expand rows if needed (up to 2x)
        if c_rows > t_rows and c_rows <= t_rows * 2:
            _add_table_rows(tmpl, c_rows - t_rows)
            t_rows = c_rows

        # Fill cells preserving formatting
        for ri in range(min(c_rows, t_rows)):
            for ci in range(min(c_cols, t_cols)):
                try:
                    _inject_table_cell_text(tmpl.cell(ri, ci), c_data[ri][ci])
                except Exception:
                    pass
        # Clear extra template cells
        for ri in range(c_rows, t_rows):
            for ci in range(t_cols):
                try:
                    _inject_table_cell_text(tmpl.cell(ri, ci), "")
                except Exception:
                    pass
    else:
        try:
            spTree = cloned_slide.shapes._spTree
            spTree.append(deepcopy(ct["element"]))
        except Exception:
            pass


def _handle_charts(
    cloned_slide, content_data: ContentData,
    src_slide, dst_prs: Presentation,
) -> None:
    """Best-effort chart transfer from content slide."""
    if not content_data.charts:
        return
    for chart_info in content_data.charts:
        try:
            chart_part = chart_info["chart_part"]
            # Copy chart part into destination
            chart_rel_type = (
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart"
            )
            new_rid = cloned_slide.part.rels.get_or_add(chart_rel_type, chart_part)
            # Add the chart graphicFrame element
            el = deepcopy(chart_info["element"])
            # Update rId in the element
            for node in el.iter():
                for attr in list(node.attrib.keys()):
                    if node.attrib[attr] in ("rId1", "rId2", "rId3"):
                        node.attrib[attr] = new_rid
            cloned_slide.shapes._spTree.append(el)
        except Exception as exc:
            log.warning("  Chart transfer failed: %s", exc)


def _handle_images(
    cloned_slide, content_data: ContentData,
    slide_w: int, slide_h: int,
) -> None:
    if not content_data.images:
        return
    occupied = []
    for shape in cloned_slide.shapes:
        occupied.append((shape.top or 0) + (shape.height or 0))

    for img in content_data.images:
        from pptx_template_transfer.models import ImageData
        if isinstance(img, ImageData):
            blob, orig_w, orig_h = img.blob, img.width, img.height
        else:
            blob, orig_w, orig_h = img[0], img[1], img[2]
        max_bottom = max(occupied, default=int(slide_h * 0.3))
        avail_top = min(max_bottom + int(Pt(10).emu), int(slide_h * 0.85))
        avail_h = slide_h - avail_top
        if avail_h < int(slide_h * 0.1):
            continue

        tw = min(orig_w, int(slide_w * 0.6))
        th_ = min(orig_h, avail_h)
        if orig_w > 0 and orig_h > 0:
            scale = min(tw / orig_w, th_ / orig_h)
            tw, th_ = int(orig_w * scale), int(orig_h * scale)

        try:
            cloned_slide.shapes.add_picture(
                io.BytesIO(blob), (slide_w - tw) // 2, avail_top, tw, th_,
            )
        except Exception:
            pass


# ============================================================================
# POST-PROCESSING
# ============================================================================

def _post_process(output_prs: Presentation) -> None:
    sw = output_prs.slide_width or 1
    sh = output_prs.slide_height or 1

    for slide_idx, slide in enumerate(output_prs.slides):
        slide_num = slide_idx + 1

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text

            # Page numbers - text pattern
            if PAGE_NUM_PATTERN.search(text):
                new = PAGE_NUM_PATTERN.sub(f"Page {slide_num:02d}", text)
                if new != text:
                    _inject_text_simple(shape, new)
                continue

            # XML-level slide number placeholder
            if has_placeholder_type(shape, {PH_SLIDE_NUMBER}):
                _inject_text_simple(shape, str(slide_num))
                continue

            # Date placeholder
            if has_placeholder_type(shape, {PH_DATE}):
                _inject_text_simple(shape, date.today().strftime("%Y-%m-%d"))
                continue

            # Dates in footer area
            if shape_bottom_frac(shape, sh) >= 0.90:
                m = DATE_PATTERN.search(text)
                if m:
                    new = DATE_PATTERN.sub(date.today().strftime("%Y-%m-%d"), text)
                    if new != text:
                        _inject_text_simple(shape, new)


def _cleanup_broken_rels(output_prs: Presentation) -> int:
    """Remove broken relationship references that prevent LibreOffice from opening.

    Returns count of removed relationships.
    """
    removed = 0
    for slide in output_prs.slides:
        part = slide.part
        bad_keys: list[str] = []
        for rel_key, rel in part.rels.items():
            try:
                if not rel.is_external:
                    _ = rel.target_part  # will throw if broken
            except Exception:
                bad_keys.append(rel_key)
        for key in bad_keys:
            try:
                del part.rels[key]
                removed += 1
            except Exception:
                pass
    return removed


def _transfer_notes(src_content: ContentData, dst_slide) -> None:
    """Copy speaker notes from content data to the output slide."""
    if not src_content.notes:
        return
    try:
        notes_slide = dst_slide.notes_slide
        tf = notes_slide.notes_text_frame
        if tf:
            existing = tf.text.strip()
            if existing:
                tf.text = existing + "\n\n---\n\n" + src_content.notes
            else:
                tf.text = src_content.notes
    except Exception:
        pass


def _validate_output(output_prs: Presentation) -> list[str]:
    """Validate the output presentation. Returns list of warnings."""
    warnings = []
    for i, slide in enumerate(output_prs.slides):
        shapes = list(slide.shapes)
        if not shapes:
            warnings.append(f"Slide {i+1}: no shapes")
    return warnings


# ============================================================================
# DIAGNOSTICS & REPORTING
# ============================================================================

def _print_slide_diagnostic(
    slide_idx: int, total: int, content_data: ContentData,
    template_idx: int, template_struct: str,
    match_score: float, injection_diag: dict,
) -> None:
    print(f"\nSlide {slide_idx+1}/{total}:")
    print(f"  Content type: {content_data.slide_type} "
          f"({content_data.word_count} words, "
          f"{len(content_data.tables)} table(s), "
          f"{len(content_data.images)} image(s))")
    print(f"  Template match: slide {template_idx+1} "
          f"(score={match_score:.0f}, type={template_struct})")

    if injection_diag.get("shapes"):
        print("  Shape classifications:")
        for s in injection_diag["shapes"]:
            preview = f' "{s["text_preview"]}"' if s["text_preview"] else ""
            print(f'    Shape "{s["name"]}" ({s["area_pct"]}% area, '
                  f'top {s["top_pct"]:.0f}%, '
                  f'conf={s["confidence"]}){preview} -> {s["role"]}')

    if injection_diag.get("injected_title"):
        print(f'  Injected: title="{injection_diag["injected_title"]}"')
    if injection_diag.get("injected_body"):
        print(f'  Injected: body ({injection_diag["injected_body"]})')
    print(f'  Cleared: {injection_diag.get("cleared_count", 0)} template text shapes')
    print(f'  Protected: {injection_diag.get("protected_count", 0)} shapes untouched')


# ============================================================================
# DESIGN MODE ORCHESTRATOR
# ============================================================================

def apply_design(
    template_path, content_path, output_path,
    config: TransferConfig,
) -> dict[str, Any]:
    """Design mode: clone template slides, inject content. Returns report dict."""
    report: dict[str, Any] = {"mode": "design", "slides": [], "warnings": [], "errors": []}

    print("\n[design] Loading presentations...")
    template_prs = Presentation(str(template_path))
    content_prs = Presentation(str(content_path))
    th = config.thresholds

    sw, sh = template_prs.slide_width, template_prs.slide_height
    ct = len(content_prs.slides)
    tt = len(template_prs.slides)
    print(f"  Template: {tt} slides, {Emu(sw).inches:.1f}\" x {Emu(sh).inches:.1f}\"")
    print(f"  Content:  {ct} slides")

    # Step 1: Extract content
    print("\n[design] Extracting content structure...")
    content_data_list = []
    for i, slide in enumerate(content_prs.slides):
        cd = extract_content(slide, i, ct, content_prs.slide_width, content_prs.slide_height, th)
        content_data_list.append(cd)
        log.debug("  Slide %d: type=%s, words=%d, title=%r, paras=%d, tables=%d, images=%d",
                   i + 1, cd.slide_type, cd.word_count, cd.title[:40],
                   len(cd.body_paragraphs), len(cd.tables), len(cd.images))

    # Step 2: Build mapping
    print("\n[design] Mapping content slides to template slides...")
    if config.slide_map:
        mapping = [config.slide_map.get(str(i + 1), 1) - 1 for i in range(ct)]
        print("  Using manual slide mapping")
    else:
        mapping = build_slide_mapping(content_prs, template_prs, content_data_list, th)

    # Template structures for diagnostics
    t_structs = [classify_template_structure(s, sw, sh, i, tt)
                 for i, s in enumerate(template_prs.slides)]

    # Step 3: Create output
    print("\n[design] Building output presentation...")
    output_prs = Presentation(str(template_path))

    prs_element = output_prs.slides._sldIdLst
    for sldId in list(prs_element):
        rId = sldId.get(f'{{{NSMAP["r"]}}}id')
        if rId:
            try:
                output_prs.part.drop_rel(rId)
            except Exception:
                pass
        prs_element.remove(sldId)

    # Step 4: Clone and inject - with per-slide error isolation
    print("\n[design] Cloning and injecting content...")
    success_count = 0
    for ci, cd in enumerate(content_data_list):
        ti = mapping[ci]
        slide_report: dict[str, Any] = {
            "index": ci + 1, "content_type": cd.slide_type,
            "template_slide": ti + 1, "template_type": t_structs[ti],
            "title": cd.title[:80], "word_count": cd.word_count,
            "status": "ok",
        }

        try:
            src_slide = template_prs.slides[ti]
            new_slide = _clone_slide(template_prs, src_slide, output_prs)

            diag = inject_content(new_slide, cd, sw, sh, th)

            _handle_tables(new_slide, cd, sw, sh)
            _handle_charts(new_slide, cd, src_slide, output_prs)
            _handle_images(new_slide, cd, sw, sh)

            if config.preserve_notes:
                _transfer_notes(cd, new_slide)

            slide_report["classifications"] = diag.get("shapes", [])
            slide_report["protected_shapes"] = diag.get("protected_count", 0)
            success_count += 1

            # Print progress
            title_preview = cd.title[:50] if cd.title else "(no title)"
            if config.verbose:
                tinfo = {
                    "struct": t_structs[ti],
                    "words": sum(word_count(text_of(s)) for s in src_slide.shapes),
                    "has_table": any(is_table(s) for s in src_slide.shapes),
                    "is_list": t_structs[ti] in ("list", "grid"),
                }
                score = _match_score(
                    cd.slide_type, tinfo["struct"], ci, ti, ct, tt,
                    cd.word_count, tinfo["words"],
                    len(cd.tables) > 0, tinfo["has_table"],
                    len(cd.body_paragraphs), tinfo["is_list"],
                )
                slide_report["match_score"] = round(score, 1)
                _print_slide_diagnostic(ci, ct, cd, ti, t_structs[ti], score, diag)
            else:
                print(f"  Slide {ci+1}/{ct}: [{cd.slide_type}] "
                      f'"{title_preview}" <- template {ti+1} ({t_structs[ti]})')

        except Exception as exc:
            slide_report["status"] = "error"
            slide_report["error"] = str(exc)
            report["errors"].append(f"Slide {ci+1}: {exc}")
            log.error("Slide %d failed: %s\n%s", ci + 1, exc, traceback.format_exc())
            print(f"  Slide {ci+1}/{ct}: ERROR - {exc}")

            # Insert blank template slide as placeholder
            try:
                fallback_layout = output_prs.slide_layouts[0]
                output_prs.slides.add_slide(fallback_layout)
            except Exception:
                pass

        report["slides"].append(slide_report)

    # Step 5: Post-processing
    print("\n[design] Post-processing...")
    _post_process(output_prs)

    # Step 5b: Clean up broken relationships (LibreOffice compatibility)
    removed_rels = _cleanup_broken_rels(output_prs)
    if removed_rels:
        log.info("  Cleaned up %d broken relationship(s)", removed_rels)

    # Step 6: Validate
    warnings = _validate_output(output_prs)
    report["warnings"] = warnings
    for w in warnings:
        log.warning(w)

    # Step 7: Save
    print(f"\n[design] Saving to {output_path}...")
    output_prs.save(str(output_path))
    print(f"[design] Done! {success_count}/{ct} slides created successfully.")
    if report["errors"]:
        print(f"  {len(report['errors'])} slide(s) had errors - see log for details.")

    return report
