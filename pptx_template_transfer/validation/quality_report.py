"""Generate per-slide quality reports for output presentations."""
from __future__ import annotations

from collections import Counter

from pptx import Presentation

from pptx_template_transfer.models import (
    ContentData, QualityReport, SlideQuality, TransferConfig,
)
from pptx_template_transfer.validation.overlap_checker import check_overlaps
from pptx_template_transfer.validation.bounds_checker import check_bounds
from pptx_template_transfer.helpers import text_of, word_count


# ---------------------------------------------------------------------------
# Duplicate / leakage detection
# ---------------------------------------------------------------------------

def _detect_text_leakage(
    output_prs: Presentation,
    content_data_list: list[ContentData] | None = None,
    min_words: int = 6,
    max_allowed_dupes: int = 1,
) -> list[str]:
    """Find non-boilerplate sentences that appear on more slides than allowed.

    Returns a list of warning strings.  Only considers sentences with
    *min_words* or more words so that short labels ("Page 01") don't
    trigger false positives.

    Tolerates repetition that originates from agenda/TOC slides (where
    titles from other slides are naturally listed as bullet points).
    """
    from pptx_template_transfer.helpers import FOOTER_PATTERNS

    # Identify agenda/TOC slides — their content naturally echoes other slides
    agenda_slide_indices: set[int] = set()
    if content_data_list:
        for si, cd in enumerate(content_data_list):
            if cd.slide_type in ("agenda", "toc"):
                agenda_slide_indices.add(si)

    # sentence -> set of slide indices where it appears
    sentence_slides: dict[str, set[int]] = {}
    for si, slide in enumerate(output_prs.slides):
        seen_on_slide: set[str] = set()
        for shape in slide.shapes:
            t = text_of(shape)
            if not t:
                continue
            for line in t.split("\n"):
                line = line.strip()
                if not line or word_count(line) < min_words:
                    continue
                if FOOTER_PATTERNS.match(line):
                    continue
                if line not in seen_on_slide:
                    seen_on_slide.add(line)
                    sentence_slides.setdefault(line, set()).add(si)

    warnings: list[str] = []
    for sentence, slides in sentence_slides.items():
        if len(slides) > max_allowed_dupes:
            # If all but one occurrence is on an agenda slide, skip it
            non_agenda = slides - agenda_slide_indices
            if len(non_agenda) <= max_allowed_dupes:
                continue
            slide_nums = sorted(s + 1 for s in slides)
            warnings.append(
                f"Text leakage: \"{sentence[:80]}\" repeated on slides "
                f"{slide_nums} ({len(slides)} occurrences)"
            )
    return warnings


def _detect_body_in_forbidden_zones(
    output_prs: Presentation,
    content_data_list: list[ContentData],
) -> list[str]:
    """Detect body content text that ended up in header/footer zones."""
    warnings: list[str] = []
    sh = output_prs.slide_height

    # Collect all unique body sentences from source content
    source_body_sentences: set[str] = set()
    for cd in content_data_list:
        for p in cd.body_paragraphs:
            if word_count(p.text) >= 5:
                source_body_sentences.add(p.text.strip())

    for si, slide in enumerate(output_prs.slides):
        for shape in slide.shapes:
            t = text_of(shape)
            if not t:
                continue
            bottom_frac = ((shape.top or 0) + (shape.height or 0)) / sh if sh else 0
            top_frac = (shape.top or 0) / sh if sh else 0

            # Check if shape is in a forbidden zone
            in_footer = bottom_frac > 0.90
            in_header = top_frac < 0.06

            if not in_footer and not in_header:
                continue

            for line in t.split("\n"):
                line = line.strip()
                if line in source_body_sentences:
                    zone = "footer" if in_footer else "header"
                    warnings.append(
                        f"Slide {si+1}: body text in {zone} zone: "
                        f"\"{line[:60]}\""
                    )
    return warnings


# ---------------------------------------------------------------------------
# Main quality report generator
# ---------------------------------------------------------------------------

def generate_quality_report(
    output_prs: Presentation,
    content_data_list: list[ContentData],
    config: TransferConfig | None = None,
) -> QualityReport:
    """Generate a comprehensive quality report for the output deck."""
    report = QualityReport()
    sw, sh = output_prs.slide_width, output_prs.slide_height

    for si, slide in enumerate(output_prs.slides):
        sq = SlideQuality(slide_index=si, build_method="native")

        # Content coverage — compare output text to source content
        if si < len(content_data_list):
            cd = content_data_list[si]
            source_words = cd.word_count
            output_text = " ".join(text_of(s) for s in slide.shapes if text_of(s))
            output_words = word_count(output_text)
            if source_words > 0:
                sq.content_coverage_pct = min(100.0, round(output_words / source_words * 100, 1))
            else:
                sq.content_coverage_pct = 100.0

            if sq.content_coverage_pct < 70:
                sq.needs_manual_review = True
                sq.review_reasons.append(
                    f"Content coverage {sq.content_coverage_pct}% (target: 70%+)"
                )

        # Font size warnings
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt < 7:
                        sq.font_warnings.append(
                            f"Shape '{shape.name}': {run.font.size.pt}pt is below minimum"
                        )
                        sq.needs_manual_review = True
                        sq.review_reasons.append("Font size below 7pt")
                        break

        # Overlap check
        sq.overlap_issues = check_overlaps(slide, si)
        if any(o.severity in ("major", "complete") for o in sq.overlap_issues):
            sq.needs_manual_review = True
            sq.review_reasons.append("Major shape overlap detected")

        # Bounds check
        sq.bounds_issues = check_bounds(slide, si, sw, sh)
        if sq.bounds_issues:
            sq.needs_manual_review = True
            sq.review_reasons.append("Shape extends beyond slide edge")

        # Empty slide
        shapes = list(slide.shapes)
        text_shapes = [s for s in shapes if text_of(s)]
        if not text_shapes:
            sq.needs_manual_review = True
            sq.review_reasons.append("Slide has no text content")

        report.slides.append(sq)

    # Per-slide composite score
    for sq in report.slides:
        slide_score = sq.content_coverage_pct
        slide_score -= len(sq.overlap_issues) * 5
        slide_score -= sum(
            10 if o.severity == "complete" else 5 if o.severity == "major" else 2
            for o in sq.overlap_issues
        )
        slide_score -= len(sq.bounds_issues) * 5
        slide_score -= len(sq.font_warnings) * 3
        sq.slide_score = max(0.0, min(100.0, slide_score))  # type: ignore[attr-defined]

        # Acceptance gate: flag slides below 50
        if sq.slide_score < 50 and not sq.needs_manual_review:  # type: ignore[attr-defined]
            sq.needs_manual_review = True
            sq.review_reasons.append(
                f"Slide quality score {sq.slide_score:.0f}/100 below acceptance threshold"  # type: ignore[attr-defined]
            )

    # Overall metrics
    report.native_count = sum(1 for s in report.slides if s.build_method == "native")
    report.fallback_count = sum(1 for s in report.slides if s.build_method == "fallback")
    report.split_count = sum(1 for s in report.slides if s.build_method == "split")

    # Overall score — weighted average of per-slide scores
    slide_scores = [getattr(s, "slide_score", s.content_coverage_pct) for s in report.slides]
    avg_score = sum(slide_scores) / len(slide_scores) if slide_scores else 100
    report.overall_score = max(0, min(100, avg_score))

    # Collect per-slide warnings
    for sq in report.slides:
        for reason in sq.review_reasons:
            report.warnings.append(f"Slide {sq.slide_index + 1}: {reason}")

    # Cross-slide validations
    leakage_warnings = _detect_text_leakage(output_prs, content_data_list)
    if leakage_warnings:
        report.warnings.extend(leakage_warnings)
        # Penalise score for leakage
        report.overall_score = max(0, report.overall_score - len(leakage_warnings) * 5)

    forbidden_zone_warnings = _detect_body_in_forbidden_zones(
        output_prs, content_data_list,
    )
    if forbidden_zone_warnings:
        report.warnings.extend(forbidden_zone_warnings)
        report.overall_score = max(0, report.overall_score - len(forbidden_zone_warnings) * 3)

    return report
