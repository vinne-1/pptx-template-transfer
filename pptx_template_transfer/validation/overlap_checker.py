"""Detect overlapping shapes on slides."""
from __future__ import annotations

from pptx_template_transfer.models import OverlapIssue


def _aabb_overlap_pct(
    ax: int, ay: int, aw: int, ah: int,
    bx: int, by: int, bw: int, bh: int,
) -> float:
    """Calculate overlap percentage between two axis-aligned bounding boxes."""
    x_overlap = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    y_overlap = max(0, min(ay + ah, by + bh) - max(ay, by))
    overlap_area = x_overlap * y_overlap
    if overlap_area == 0:
        return 0.0
    smaller_area = min(aw * ah, bw * bh)
    if smaller_area == 0:
        return 0.0
    return overlap_area / smaller_area * 100.0


def _is_decorative(shape) -> bool:
    """Return True if the shape is a non-text decorative element (accent,
    background shape, line, etc.) that should not trigger overlap warnings."""
    has_text = (
        hasattr(shape, 'has_text_frame')
        and shape.has_text_frame
        and shape.text_frame.text.strip()
    )
    if has_text:
        return False
    # Auto-shapes without text are decorative (ovals, rectangles, triangles, lines)
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type in (
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
            MSO_SHAPE_TYPE.LINE,
            1,   # RECTANGLE enum
            9,   # OVAL enum
        ):
            return True
    except Exception:
        pass
    return not has_text


def check_overlaps(slide, slide_index: int) -> list[OverlapIssue]:
    """Check all shape pairs on a slide for overlaps."""
    issues: list[OverlapIssue] = []
    shapes = list(slide.shapes)

    for i in range(len(shapes)):
        for j in range(i + 1, len(shapes)):
            a, b = shapes[i], shapes[j]
            ax, ay = a.left or 0, a.top or 0
            aw, ah = a.width or 0, a.height or 0
            bx, by = b.left or 0, b.top or 0
            bw, bh = b.width or 0, b.height or 0

            if aw == 0 or ah == 0 or bw == 0 or bh == 0:
                continue

            pct = _aabb_overlap_pct(ax, ay, aw, ah, bx, by, bw, bh)
            if pct < 5:
                continue

            severity = "minor" if pct < 25 else ("major" if pct < 75 else "complete")

            # Skip overlaps involving decorative shapes (background accents)
            a_is_text = hasattr(a, 'has_text_frame') and a.has_text_frame and a.text_frame.text.strip()
            b_is_text = hasattr(b, 'has_text_frame') and b.has_text_frame and b.text_frame.text.strip()
            if not a_is_text and not b_is_text:
                continue  # Both decorative — overlap is fine
            if _is_decorative(a) or _is_decorative(b):
                continue  # One is decorative — background accent overlapping content is ok

            issues.append(OverlapIssue(
                slide_index=slide_index,
                shape_a=a.name or f"shape_{i}",
                shape_b=b.name or f"shape_{j}",
                overlap_pct=round(pct, 1),
                severity=severity,
            ))

    return issues
