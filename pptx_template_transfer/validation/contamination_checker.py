"""Detect target-deck body text contamination in the output."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from pptx_template_transfer.helpers import (
    FOOTER_PATTERNS, text_of, word_count,
)


def _extract_body_ngrams(
    text: str, n: int = 3, min_words: int = 4,
) -> set[tuple[str, ...]]:
    """Extract word n-grams from text, skipping short fragments."""
    words = text.lower().split()
    if len(words) < min_words:
        return set()
    return {tuple(words[i:i + n]) for i in range(len(words) - n + 1)}


def _slide_body_text(slide, slide_height: int = 0) -> str:
    """Extract all body-zone text from a slide (excluding footer/header)."""
    parts: list[str] = []
    for shape in slide.shapes:
        t = text_of(shape)
        if not t or word_count(t) < 3:
            continue
        if FOOTER_PATTERNS.match(t.strip()):
            continue
        # Skip shapes in bottom 10% (footer zone)
        if slide_height:
            bottom_frac = ((shape.top or 0) + (shape.height or 0)) / slide_height
            if bottom_frac > 0.90 and word_count(t) <= 15:
                continue
        parts.append(t)
    return " ".join(parts)


def check_target_contamination(
    output_prs: Presentation,
    target_path: Path,
    similarity_threshold: float = 0.40,
) -> list[str]:
    """Compare output body text against target deck body text.

    For each output slide, compute Jaccard similarity of word 3-grams
    against every target slide.  If similarity exceeds the threshold,
    flag it as target-content contamination.

    Returns a list of warning strings.
    """
    target_prs = Presentation(str(target_path))
    target_sh = target_prs.slide_height
    output_sh = output_prs.slide_height

    # Pre-compute target slide n-grams
    target_ngrams: list[tuple[int, set[tuple[str, ...]]]] = []
    for ti, tslide in enumerate(target_prs.slides):
        text = _slide_body_text(tslide, target_sh)
        ngrams = _extract_body_ngrams(text)
        if ngrams:
            target_ngrams.append((ti, ngrams))

    warnings: list[str] = []
    for oi, oslide in enumerate(output_prs.slides):
        otext = _slide_body_text(oslide, output_sh)
        ongrams = _extract_body_ngrams(otext)
        if not ongrams:
            continue

        for ti, tngrams in target_ngrams:
            if not tngrams:
                continue
            intersection = ongrams & tngrams
            union = ongrams | tngrams
            jaccard = len(intersection) / len(union) if union else 0.0

            if jaccard > similarity_threshold:
                warnings.append(
                    f"Output slide {oi+1} has {jaccard:.0%} similarity with "
                    f"target slide {ti+1} — possible target-content contamination"
                )

    return warnings
