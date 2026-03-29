"""Compute source-deck content coverage in the output."""
from __future__ import annotations

from pptx import Presentation

from pptx_template_transfer.helpers import text_of, word_count
from pptx_template_transfer.models import (
    ContentData,
    SourceCoverageEntry,
    SourceCoverageReport,
)


def _extract_output_words_per_slide(output_prs: Presentation) -> list[set[str]]:
    """Extract lowercased word sets for each output slide."""
    result: list[set[str]] = []
    for slide in output_prs.slides:
        words: set[str] = set()
        for shape in slide.shapes:
            t = text_of(shape)
            if t:
                words.update(t.lower().split())
        result.append(words)
    return result


def compute_source_coverage(
    output_prs: Presentation,
    content_data_list: list[ContentData],
) -> SourceCoverageReport:
    """For every source slide, compute how much of its content made it
    into the output.

    Uses word-level overlap: for each source slide's text, check what
    fraction of its words appear in the mapped output slide(s).
    """
    report = SourceCoverageReport(
        total_source_slides=len(content_data_list),
        total_output_slides=len(output_prs.slides),
    )

    output_word_sets = _extract_output_words_per_slide(output_prs)

    total_source_words = 0
    total_covered_words = 0

    for si, cd in enumerate(content_data_list):
        entry = SourceCoverageEntry(
            source_slide_index=si,
            source_title=cd.title[:80] if cd.title else "",
            source_word_count=cd.word_count,
        )

        # Collect all source words for this slide
        source_text_parts = [cd.title] + [p.text for p in cd.body_paragraphs]
        source_words = set()
        for part in source_text_parts:
            if part:
                source_words.update(part.lower().split())

        # Find best-matching output slide(s)
        # The primary mapping is 1:1 by index (source slide i → output slide i)
        # but we also check adjacent slides in case of splits
        best_coverage = 0.0
        mapped_indices: list[int] = []

        search_range = range(
            max(0, si - 1),
            min(len(output_word_sets), si + 3),
        )
        for oi in search_range:
            if not source_words:
                break
            overlap = source_words & output_word_sets[oi]
            coverage = len(overlap) / len(source_words) if source_words else 0.0
            if coverage > 0.15:  # at least 15% overlap to count as mapped
                mapped_indices.append(oi)
                best_coverage = max(best_coverage, coverage)

        # Also check the direct 1:1 mapping
        if si < len(output_word_sets) and si not in mapped_indices:
            overlap = source_words & output_word_sets[si]
            coverage = len(overlap) / len(source_words) if source_words else 0.0
            if coverage > 0.05:
                mapped_indices.append(si)
                best_coverage = max(best_coverage, coverage)

        entry.output_slide_indices = sorted(set(mapped_indices))
        entry.text_used_pct = round(best_coverage * 100, 1)

        # Block-level coverage: check each text_block individually
        all_output_words: set[str] = set()
        for oi in mapped_indices:
            all_output_words |= output_word_sets[oi]

        blocks = cd.text_blocks or []
        entry.blocks_total = len(blocks)
        entry.blocks_covered = 0
        entry.missing_block_texts = []
        for block in blocks:
            block_words = set()
            for para in block.paragraphs:
                if para.text:
                    block_words.update(para.text.lower().split())
            if not block_words:
                entry.blocks_covered += 1  # empty block counts as covered
                continue
            overlap = block_words & all_output_words
            block_cov = len(overlap) / len(block_words)
            if block_cov >= 0.40:
                entry.blocks_covered += 1
            else:
                # Store first 120 chars of block text for diagnostics
                block_text = " ".join(p.text for p in block.paragraphs)
                entry.missing_block_texts.append(block_text[:120])

        # Track dropped assets
        entry.tables_dropped = len(cd.tables)  # TODO: detect if tables were rebuilt
        entry.charts_dropped = len(cd.charts)
        entry.images_dropped = len(cd.images)

        if not mapped_indices:
            report.unmapped_source_slides.append(si)
            report.warnings.append(
                f"Source slide {si+1} '{cd.title[:50]}' has no matching output slide"
            )

        if source_words and best_coverage < 0.30:
            report.warnings.append(
                f"Source slide {si+1} '{cd.title[:50]}': only {entry.text_used_pct}% "
                f"text coverage"
            )

        total_source_words += len(source_words)
        total_covered_words += int(len(source_words) * best_coverage)

        report.entries.append(entry)

    report.overall_text_coverage_pct = round(
        total_covered_words / total_source_words * 100 if total_source_words else 100.0,
        1,
    )

    if report.overall_text_coverage_pct < 70:
        report.warnings.insert(0,
            f"Overall source content coverage is {report.overall_text_coverage_pct}% "
            f"(target: 70%+)"
        )

    return report
