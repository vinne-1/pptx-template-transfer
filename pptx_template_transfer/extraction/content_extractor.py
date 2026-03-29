"""Extract structured content from each source slide."""
from __future__ import annotations

import re
from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation

from pptx_template_transfer.models import (
    ContentData, ParagraphData, RunData, TextBlock, Thresholds,
)
from pptx_template_transfer.helpers import (
    FOOTER_PATTERNS, PH_FOOTER_SET, PH_TITLE_SET,
    text_of, word_count, max_font_pt, shape_area_pct,
    is_picture, is_table, is_chart, is_group, is_ole_or_embedded,
    dominant_text_color, placeholder_type_int,
)
from pptx_template_transfer.analysis.slide_classifier import classify_slide_type


def _extract_paragraphs_from_shape(shape) -> list[ParagraphData]:
    result = []
    if not shape.has_text_frame:
        return result
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level = para.level if para.level else 0
        bold = italic = False
        font_size = 0.0
        runs_data = []
        for run in para.runs:
            r_bold = bool(run.font.bold)
            r_italic = bool(run.font.italic)
            r_size = run.font.size.pt if run.font.size else 0.0
            r_url = None
            try:
                if run.hyperlink and run.hyperlink.address:
                    r_url = run.hyperlink.address
            except Exception:
                pass
            runs_data.append(RunData(
                text=run.text, bold=r_bold, italic=r_italic,
                font_size=r_size, hyperlink_url=r_url,
            ))
            if r_bold:
                bold = True
            if r_italic:
                italic = True
            font_size = max(font_size, r_size)
        result.append(ParagraphData(
            text=text, level=level, bold=bold,
            italic=italic, font_size=font_size, runs=runs_data,
        ))
    return result


def _extract_table_data(shape) -> list[list[str]]:
    if not is_table(shape):
        return []
    rows = []
    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def _extract_chart_info(shape, slide) -> dict | None:
    if not is_chart(shape):
        return None
    try:
        chart_part = shape.chart_part
        return {
            "element": deepcopy(shape._element),
            "chart_part": chart_part,
            "width": shape.width,
            "height": shape.height,
            "left": shape.left,
            "top": shape.top,
        }
    except Exception:
        return None


def extract_content(
    slide, slide_index: int, total_slides: int,
    slide_w: int, slide_h: int, th: Thresholds,
) -> ContentData:
    """Extract structured content from a single slide."""
    content = ContentData()
    content.source_slide_index = slide_index
    content.slide_type = classify_slide_type(slide, slide_index, total_slides, slide_w, slide_h)
    shapes = list(slide.shapes)

    # --- Title detection ---
    # Score each text shape as a title candidate using font size, vertical
    # position, word count, and whether it sits in a footer/boilerplate zone.
    text_shapes = [(s, max_font_pt(s), text_of(s)) for s in shapes if text_of(s)]

    def _title_score(item: tuple) -> float:
        s, fs, txt = item
        wc = word_count(txt)
        top_frac = (s.top or 0) / slide_h if slide_h else 0.0
        bottom_frac = ((s.top or 0) + (s.height or 0)) / slide_h if slide_h else 1.0

        # Disqualify footer-zone shapes and very long text
        if bottom_frac > 0.88 and wc <= 12:
            return -1.0
        if wc > 20:
            return -1.0
        if FOOTER_PATTERNS.match(txt.strip()):
            return -1.0
        if placeholder_type_int(s) is not None and placeholder_type_int(s) in PH_FOOTER_SET:
            return -1.0

        score = 0.0
        # Font size bonus (dominant signal)
        score += min(fs, 40.0) * 2.0
        # Position: top-of-slide titles score higher
        score += max(0.0, (1.0 - top_frac)) * 20.0
        # Prefer concise titles
        if wc <= 8:
            score += 10.0
        elif wc <= 15:
            score += 5.0
        # Placeholder type bonus
        pt = placeholder_type_int(s)
        if pt is not None and pt in PH_TITLE_SET:
            score += 30.0
        # Bold emphasis bonus
        if s.has_text_frame:
            first_runs = s.text_frame.paragraphs[0].runs if s.text_frame.paragraphs else []
            if first_runs and any(r.font.bold for r in first_runs):
                score += 5.0
        return score

    title_shape = None
    if text_shapes:
        scored = [(s, fs, txt, _title_score((s, fs, txt))) for s, fs, txt in text_shapes]
        scored.sort(key=lambda x: -x[3])
        if scored and scored[0][3] > 0:
            title_shape = scored[0][0]
            content.title = scored[0][2]

    # --- Add title shape as a positioned TextBlock ---
    if title_shape and text_of(title_shape) and slide_w > 0 and slide_h > 0:
        paras = _extract_paragraphs_from_shape(title_shape)
        if paras:
            content.text_blocks.append(TextBlock(
                paragraphs=paras,
                left_pct=(title_shape.left or 0) / slide_w * 100,
                top_pct=(title_shape.top or 0) / slide_h * 100,
                width_pct=(title_shape.width or 0) / slide_w * 100,
                height_pct=(title_shape.height or 0) / slide_h * 100,
                is_heading=True,
                is_label=False,
            ))

    # --- Body extraction (with zone protection) ---
    body_shapes = [
        s for s in shapes
        if s is not title_shape and text_of(s) and not is_table(s) and not is_chart(s)
    ]

    def _is_footer_zone_shape(s) -> bool:
        """Return True if the shape sits in a header or footer zone and should
        be excluded from body content extraction."""
        txt = text_of(s).strip()
        # Explicit footer patterns (page numbers, confidential, dates)
        if FOOTER_PATTERNS.match(txt):
            return True
        # Footer placeholder types
        if placeholder_type_int(s) is not None and placeholder_type_int(s) in PH_FOOTER_SET:
            return True
        # Spatial: shape in the bottom 10% of the slide with short text
        bottom_frac = ((s.top or 0) + (s.height or 0)) / slide_h if slide_h else 0.0
        top_frac = (s.top or 0) / slide_h if slide_h else 0.0
        wc = word_count(txt)
        area_pct = shape_area_pct(s, slide_w, slide_h)
        if bottom_frac > 0.90 and wc <= 15 and area_pct < 8.0:
            return True
        # Header-zone: tiny text in the top 8% of the slide
        if top_frac < 0.08 and wc <= 5 and area_pct < 3.0:
            return True
        return False

    body_shapes = [s for s in body_shapes if not _is_footer_zone_shape(s)]
    body_shapes.sort(key=lambda s: ((s.top or 0), (s.left or 0)))

    for shape in body_shapes:
        paras = _extract_paragraphs_from_shape(shape)
        for p in paras:
            if p.bold or (p.font_size >= th.subheading_min_font_pt and word_count(p.text) <= 10):
                p.bold = True
            content.body_paragraphs.append(p)

        if paras and slide_w > 0 and slide_h > 0:
            text = text_of(shape)
            wc = word_count(text)
            max_fs = max_font_pt(shape)
            is_heading = (max_fs >= th.subheading_min_font_pt and wc <= 10) or (
                paras[0].bold and wc <= 10
            )
            is_label = wc <= 3 and max_fs < 20
            content.text_blocks.append(TextBlock(
                paragraphs=paras,
                left_pct=(shape.left or 0) / slide_w * 100,
                top_pct=(shape.top or 0) / slide_h * 100,
                width_pct=(shape.width or 0) / slide_w * 100,
                height_pct=(shape.height or 0) / slide_h * 100,
                is_heading=is_heading,
                is_label=is_label,
            ))

    # --- Tables ---
    for shape in shapes:
        if is_table(shape):
            table_text = _extract_table_data(shape)
            content.tables.append({
                "data": table_text,
                "rows": len(table_text),
                "cols": len(table_text[0]) if table_text else 0,
                "element": deepcopy(shape._element),
                "width": shape.width, "height": shape.height,
                "left": shape.left, "top": shape.top,
            })

    # --- Charts ---
    for shape in shapes:
        if is_chart(shape):
            content.has_chart = True
            ci = _extract_chart_info(shape, slide)
            if ci:
                content.charts.append(ci)

    # --- Images ---
    for shape in shapes:
        if is_picture(shape):
            area_pct = shape_area_pct(shape, slide_w, slide_h)
            if area_pct > th.image_min_area_pct:
                try:
                    blob = shape.image.blob
                    content.images.append(
                        (blob, shape.width, shape.height, shape.left or 0, shape.top or 0),
                    )
                except Exception:
                    pass

    # --- Speaker notes ---
    try:
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                content.notes = notes_tf.text.strip()
    except Exception:
        pass

    # --- Word count and color ---
    all_text = content.title + " " + " ".join(p.text for p in content.body_paragraphs)
    content.word_count = word_count(all_text)
    colors: dict[str, int] = {}
    for shape in shapes:
        c = dominant_text_color(shape)
        if c:
            colors[c] = colors.get(c, 0) + 1
    if colors:
        content.primary_color = max(colors, key=colors.get)

    # --- Semantic blocks ---
    from pptx_template_transfer.extraction.semantic_blocks import detect_semantic_blocks
    content.semantic_blocks = detect_semantic_blocks(content.body_paragraphs)

    return content


def extract_all_content(
    content_path: Path, th: Thresholds,
) -> list[ContentData]:
    """Extract structured content from every slide in a PPTX."""
    prs = Presentation(str(content_path))
    sw, sh = prs.slide_width, prs.slide_height
    ct = len(prs.slides)
    result = []
    for i, slide in enumerate(prs.slides):
        cd = extract_content(slide, i, ct, sw, sh, th)
        result.append(cd)
    return result
