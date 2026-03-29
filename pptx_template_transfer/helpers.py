"""Shared utility functions and constants for PPTX Template Transfer."""
from __future__ import annotations

import re
from typing import Any

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

# ============================================================================
# XML NAMESPACES
# ============================================================================

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
}

# ============================================================================
# REGEX PATTERNS
# ============================================================================

FOOTER_PATTERNS = re.compile(
    r"(?i)(page\s*\d+|confidential|©|\bcopyright\b|\ball rights reserved\b"
    r"|\b\d{4}[-/]\d{2}[-/]\d{2}\b|\b\d{2}/\d{2}/\d{4}\b"
    r"|proprietary|internal use|draft|do not distribute)",
)
PAGE_NUM_PATTERN = re.compile(r"(?i)page\s*\d+")
DATE_PATTERN = re.compile(
    r"\b(\d{4}[-/]\d{1,2}[-/]\d{1,2}|\d{1,2}/\d{1,2}/\d{4})\b",
)
JUST_NUMBER_RE = re.compile(r"^\d{1,2}$")

# ============================================================================
# PLACEHOLDER CONSTANTS
# ============================================================================

PH_TITLE = 15
PH_CENTER_TITLE = 3
PH_SUBTITLE = 4
PH_BODY = 2
PH_OBJECT = 7
PH_FOOTER = 11
PH_SLIDE_NUMBER = 12
PH_DATE = 10

PH_TITLE_SET = {PH_TITLE, PH_CENTER_TITLE}
PH_BODY_SET = {PH_BODY, PH_OBJECT, PH_SUBTITLE}
PH_FOOTER_SET = {PH_FOOTER, PH_SLIDE_NUMBER, PH_DATE}


# ============================================================================
# SHAPE HELPERS
# ============================================================================

def text_of(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def word_count(text: str) -> int:
    return len(text.split()) if text else 0


def max_font_pt(shape) -> float:
    mx = 0.0
    if not shape.has_text_frame:
        return mx
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                mx = max(mx, run.font.size.pt)
    return mx


def min_font_pt(shape) -> float:
    mn = 999.0
    if not shape.has_text_frame:
        return 0.0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None and run.font.size.pt > 0:
                mn = min(mn, run.font.size.pt)
    return mn if mn < 999.0 else 0.0


def shape_area(shape) -> int:
    return (shape.width or 0) * (shape.height or 0)


def shape_area_pct(shape, slide_w: int, slide_h: int) -> float:
    total = slide_w * slide_h
    return shape_area(shape) / total * 100.0 if total else 0.0


def is_picture(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13
    except Exception:
        return False


def is_table(shape) -> bool:
    return hasattr(shape, "has_table") and shape.has_table


def is_chart(shape) -> bool:
    return hasattr(shape, "has_chart") and shape.has_chart


def is_group(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def is_ole_or_embedded(shape) -> bool:
    try:
        tag = etree.QName(shape._element.tag).localname
        return tag == "graphicFrame" and not is_table(shape) and not is_chart(shape)
    except Exception:
        return False


def shape_bottom_frac(shape, slide_h: int) -> float:
    if slide_h == 0:
        return 0.0
    return ((shape.top or 0) + (shape.height or 0)) / slide_h


def shape_top_frac(shape, slide_h: int) -> float:
    return (shape.top or 0) / slide_h if slide_h else 0.0


def shape_left_frac(shape, slide_w: int) -> float:
    return (shape.left or 0) / slide_w if slide_w else 0.0


def is_allcaps_short(text: str) -> bool:
    words = text.split()
    if not words or len(words) > 5:
        return False
    alpha = "".join(c for c in text if c.isalpha())
    return bool(alpha) and alpha == alpha.upper()


def dominant_text_color(shape) -> str | None:
    if not shape.has_text_frame:
        return None
    colors: dict[str, int] = {}
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                c = run.font.color
                if c and c.type is not None and c.rgb:
                    key = str(c.rgb)
                    colors[key] = colors.get(key, 0) + len(run.text)
            except (AttributeError, TypeError):
                pass
    return max(colors, key=colors.get) if colors else None


def group_text_words(shape) -> int:
    total = 0
    try:
        for child in shape.shapes:
            total += word_count(text_of(child))
            if is_group(child):
                total += group_text_words(child)
    except Exception:
        pass
    return total


def has_placeholder_type(shape, ph_types: set[int]) -> bool:
    try:
        pf = shape.placeholder_format
        if pf is not None and pf.type is not None:
            return int(pf.type) in ph_types
    except Exception:
        pass
    return False


def placeholder_type_int(shape) -> int | None:
    try:
        pf = shape.placeholder_format
        if pf is not None and pf.type is not None:
            return int(pf.type)
    except Exception:
        pass
    return None


def rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def update_rids_in_tree(element, rid_map: dict[str, str]) -> None:
    for el in element.iter():
        for attr_name in list(el.attrib.keys()):
            val = el.attrib[attr_name]
            if val in rid_map:
                el.attrib[attr_name] = rid_map[val]


def style_runs(
    paragraph, *,
    font_name: str, font_size_pt: float,
    bold: bool = False, italic: bool = False,
    color_hex: str = "111827",
) -> None:
    """Apply font properties to every run in a paragraph."""
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color_hex)
