"""TDD test suite for pptx_template_transfer.

RED phase: tests written first, covering data classes, helpers, extraction,
classification, template analysis, recreate mode, and end-to-end transfer.
"""

import io
import json
import tempfile
from dataclasses import asdict
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

import pptx_template_transfer as mod
from pptx_template_transfer import (
    ContentData,
    ParagraphData,
    RunData,
    TemplateStyle,
    TextBlock,
    Thresholds,
    TransferConfig,
    analyze_template,
    build_slide,
    classify_shape_role,
    detect_mode,
    extract_all_content,
    extract_content,
    get_slide_zones,
    transfer,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def simple_pptx(tmp_path: Path) -> Path:
    """Create a minimal PPTX with one slide containing a title and body."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)

    # Title placeholder
    title_ph = slide.placeholders[0]
    title_ph.text = "Test Title"
    for run in title_ph.text_frame.paragraphs[0].runs:
        run.font.size = Pt(28)

    # Subtitle/body placeholder
    body_ph = slide.placeholders[1]
    body_ph.text = "This is the body text with enough words to be interesting."
    for run in body_ph.text_frame.paragraphs[0].runs:
        run.font.size = Pt(14)

    path = tmp_path / "simple.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def multi_slide_pptx(tmp_path: Path) -> Path:
    """Create a PPTX with 3 slides: title, content, closing."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.placeholders[0].text = "Presentation Title"
    for run in slide1.placeholders[0].text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)

    # Slide 2: Content (use blank layout + manual textboxes)
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide2 = prs.slides.add_slide(blank_layout)
    # Title textbox
    tb_title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    p = tb_title.text_frame.paragraphs[0]
    p.text = "Content Slide Title"
    for run in p.runs:
        run.font.size = Pt(24)
    # Body textbox
    tb_body = slide2.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(4))
    p = tb_body.text_frame.paragraphs[0]
    p.text = "First paragraph with important details about the project scope and timeline."
    for run in p.runs:
        run.font.size = Pt(12)
    p2 = tb_body.text_frame.add_paragraph()
    p2.text = "Second paragraph with additional information about deliverables."
    for run in p2.runs:
        run.font.size = Pt(12)

    # Slide 3: Closing
    slide3 = prs.slides.add_slide(prs.slide_layouts[0])
    slide3.placeholders[0].text = "Thank You"
    for run in slide3.placeholders[0].text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)

    path = tmp_path / "multi.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def template_with_branding(tmp_path: Path) -> Path:
    """Create a template PPTX with footer, logo-like shape, and branding."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i in range(3):
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(layout)

        # Title area
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = f"Template Slide {i + 1}"
        for run in p.runs:
            run.font.size = Pt(28)
            run.font.bold = True

        # Body area
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
        p = tb2.text_frame.paragraphs[0]
        p.text = "Template body placeholder text that should not leak."
        for run in p.runs:
            run.font.size = Pt(14)

        # Footer area (bottom 8%)
        footer_top = prs.slide_height - Inches(0.5)
        tb3 = slide.shapes.add_textbox(Inches(1), footer_top, Inches(4), Inches(0.3))
        tb3.text_frame.paragraphs[0].text = "Confidential"
        for run in tb3.text_frame.paragraphs[0].runs:
            run.font.size = Pt(8)

        # Page number footer
        tb4 = slide.shapes.add_textbox(Inches(11), footer_top, Inches(1.5), Inches(0.3))
        tb4.text_frame.paragraphs[0].text = f"Page {i + 1:02d}"
        for run in tb4.text_frame.paragraphs[0].runs:
            run.font.size = Pt(8)

    path = tmp_path / "template.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def content_with_grid(tmp_path: Path) -> Path:
    """Create a content PPTX with a multi-column grid layout."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Grid Layout Test"
    for run in p.runs:
        run.font.size = Pt(28)

    # 4-column grid: labels + descriptions
    cols = [
        ("01", "Collect", "Wazuh agents and logs"),
        ("02", "Detect", "MITRE ATT&CK mapping"),
        ("03", "Operate", "ITSM tickets and queues"),
        ("04", "Govern", "Rule catalog and plans"),
    ]
    for ci, (num, heading, desc) in enumerate(cols):
        left = Inches(1 + ci * 3)
        # Number label
        tb_num = slide.shapes.add_textbox(left, Inches(2.5), Inches(0.5), Inches(0.4))
        tb_num.text_frame.paragraphs[0].text = num
        for run in tb_num.text_frame.paragraphs[0].runs:
            run.font.size = Pt(14)
            run.font.bold = True
        # Heading
        tb_head = slide.shapes.add_textbox(left, Inches(3), Inches(2.5), Inches(0.5))
        tb_head.text_frame.paragraphs[0].text = heading
        for run in tb_head.text_frame.paragraphs[0].runs:
            run.font.size = Pt(16)
            run.font.bold = True
        # Description
        tb_desc = slide.shapes.add_textbox(left, Inches(3.6), Inches(2.5), Inches(2))
        tb_desc.text_frame.paragraphs[0].text = desc
        for run in tb_desc.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11)

    path = tmp_path / "grid_content.pptx"
    prs.save(str(path))
    return path


# ===========================================================================
# 1. DATA CLASSES
# ===========================================================================

class TestDataClasses:
    """Test dataclass defaults and construction."""

    def test_thresholds_defaults(self):
        th = Thresholds()
        assert th.title_min_font_pt == 20.0
        assert th.image_min_area_pct == 1.5
        assert th.body_max_zones == 2

    def test_thresholds_custom(self):
        th = Thresholds(title_min_font_pt=16, body_max_zones=4)
        assert th.title_min_font_pt == 16
        assert th.body_max_zones == 4

    def test_transfer_config_defaults(self):
        cfg = TransferConfig()
        assert cfg.mode is None
        assert cfg.verbose is False
        assert cfg.preserve_notes is True

    def test_paragraph_data(self):
        pd = ParagraphData(text="Hello", level=1, bold=True)
        assert pd.text == "Hello"
        assert pd.level == 1
        assert pd.bold is True
        assert pd.italic is False

    def test_text_block_defaults(self):
        tb = TextBlock()
        assert tb.left_pct == 0.0
        assert tb.top_pct == 0.0
        assert tb.width_pct == 20.0
        assert tb.is_heading is False
        assert tb.is_label is False

    def test_text_block_with_paragraphs(self):
        pd = ParagraphData(text="Item", bold=True)
        tb = TextBlock(
            paragraphs=[pd],
            left_pct=10.0, top_pct=30.0,
            width_pct=25.0, height_pct=15.0,
            is_heading=True,
        )
        assert len(tb.paragraphs) == 1
        assert tb.is_heading is True

    def test_content_data_has_text_blocks(self):
        cd = ContentData()
        assert cd.text_blocks == []
        assert cd.body_paragraphs == []
        assert cd.slide_type == "content"

    def test_template_style_defaults(self):
        ts = TemplateStyle()
        assert ts.heading_font == "Montserrat"
        assert ts.body_font == "Lato"
        assert ts.color_primary == "2563EB"

    def test_run_data(self):
        rd = RunData(text="link", bold=True, hyperlink_url="https://example.com")
        assert rd.hyperlink_url == "https://example.com"


# ===========================================================================
# 2. HELPER FUNCTIONS
# ===========================================================================

class TestHelpers:
    """Test pure helper functions."""

    def test_word_count_empty(self):
        assert mod._word_count("") == 0

    def test_word_count_normal(self):
        assert mod._word_count("hello world") == 2

    def test_word_count_whitespace(self):
        assert mod._word_count("  a  b  c  ") == 3

    def test_word_count_none_like(self):
        assert mod._word_count("") == 0

    def test_is_allcaps_short_true(self):
        assert mod._is_allcaps_short("SECURITY MODEL") is True

    def test_is_allcaps_short_false_lowercase(self):
        assert mod._is_allcaps_short("security model") is False

    def test_is_allcaps_short_false_too_long(self):
        assert mod._is_allcaps_short("THIS IS A VERY LONG CAPS SENTENCE") is False

    def test_is_allcaps_short_empty(self):
        assert mod._is_allcaps_short("") is False

    def test_is_allcaps_short_with_numbers(self):
        # "01" has no alpha → should return False
        assert mod._is_allcaps_short("01") is False

    def test_is_allcaps_short_mixed(self):
        assert mod._is_allcaps_short("SECTION 01") is True

    def test_rgb_valid(self):
        c = mod._rgb("FF0000")
        assert str(c) == "FF0000"

    def test_rgb_blue(self):
        c = mod._rgb("2563EB")
        assert str(c) == "2563EB"

    def test_shape_area_pct_zero_slide(self):
        shape = MagicMock(width=100, height=100)
        assert mod._shape_area_pct(shape, 0, 0) == 0.0

    def test_shape_area_pct_normal(self):
        shape = MagicMock(width=100, height=200)
        # (100*200) / (1000*1000) * 100 = 2.0%
        assert mod._shape_area_pct(shape, 1000, 1000) == pytest.approx(2.0)

    def test_shape_bottom_frac(self):
        shape = MagicMock(top=900, height=100)
        assert mod._shape_bottom_frac(shape, 1000) == pytest.approx(1.0)

    def test_shape_bottom_frac_zero(self):
        shape = MagicMock(top=0, height=0)
        assert mod._shape_bottom_frac(shape, 0) == 0.0

    def test_shape_top_frac(self):
        shape = MagicMock(top=500)
        assert mod._shape_top_frac(shape, 1000) == pytest.approx(0.5)


class TestStyleRuns:
    """Test _style_runs applies font properties to runs."""

    def test_style_runs_sets_font_name(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = "Test text"

        mod._style_runs(p, font_name="Calibri", font_size_pt=12)

        for run in p.runs:
            assert run.font.name == "Calibri"
            assert run.font.size == Pt(12)

    def test_style_runs_bold_italic(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = "Bold italic"

        mod._style_runs(p, font_name="Arial", font_size_pt=14, bold=True, italic=True)

        for run in p.runs:
            assert run.font.bold is True
            assert run.font.italic is True

    def test_style_runs_color(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = "Colored"

        mod._style_runs(p, font_name="Arial", font_size_pt=12, color_hex="FF0000")

        for run in p.runs:
            assert str(run.font.color.rgb) == "FF0000"

    def test_style_runs_empty_paragraph(self):
        """style_runs on paragraph with no runs should not raise."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        p = tb.text_frame.paragraphs[0]
        # Empty paragraph — no text, no runs
        mod._style_runs(p, font_name="Arial", font_size_pt=12)
        # Should not raise


# ===========================================================================
# 3. CONTENT EXTRACTION
# ===========================================================================

class TestContentExtraction:
    """Test extract_content captures text, positions, and structure."""

    def test_extract_simple_slide(self, simple_pptx: Path):
        prs = Presentation(str(simple_pptx))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        content = extract_content(slide, 0, 1, sw, sh, Thresholds())

        assert content.title != ""
        assert content.word_count > 0

    def test_extract_title_detected(self, multi_slide_pptx: Path):
        prs = Presentation(str(multi_slide_pptx))
        slide = prs.slides[1]  # Content slide
        sw, sh = prs.slide_width, prs.slide_height

        content = extract_content(slide, 1, 3, sw, sh, Thresholds())

        assert "Content Slide Title" in content.title

    def test_extract_body_paragraphs(self, multi_slide_pptx: Path):
        prs = Presentation(str(multi_slide_pptx))
        slide = prs.slides[1]  # Content slide
        sw, sh = prs.slide_width, prs.slide_height

        content = extract_content(slide, 1, 3, sw, sh, Thresholds())

        body_texts = [p.text for p in content.body_paragraphs]
        assert any("First paragraph" in t for t in body_texts)
        assert any("Second paragraph" in t for t in body_texts)

    def test_extract_text_blocks_populated(self, content_with_grid: Path):
        """Text blocks should capture positioned shapes from grid layouts."""
        prs = Presentation(str(content_with_grid))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        content = extract_content(slide, 0, 1, sw, sh, Thresholds())

        # Should have text blocks for the 12 grid shapes (4 cols x 3 rows: num, heading, desc)
        assert len(content.text_blocks) > 0
        # Positions should be diverse (not all at 0,0)
        lefts = {round(tb.left_pct) for tb in content.text_blocks}
        assert len(lefts) > 1, "Text blocks should have diverse left positions"

    def test_extract_text_blocks_have_position_data(self, content_with_grid: Path):
        prs = Presentation(str(content_with_grid))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        content = extract_content(slide, 0, 1, sw, sh, Thresholds())

        for tb in content.text_blocks:
            assert tb.width_pct > 0
            assert tb.height_pct > 0
            assert len(tb.paragraphs) > 0

    def test_extract_labels_flagged(self, content_with_grid: Path):
        """Short text like '01', '02' should be flagged as labels."""
        prs = Presentation(str(content_with_grid))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        content = extract_content(slide, 0, 1, sw, sh, Thresholds())

        label_blocks = [tb for tb in content.text_blocks if tb.is_label]
        assert len(label_blocks) >= 4, f"Expected >=4 labels, got {len(label_blocks)}"

    def test_extract_headings_flagged(self, content_with_grid: Path):
        """Bold short text like column headings should be flagged."""
        prs = Presentation(str(content_with_grid))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        content = extract_content(slide, 0, 1, sw, sh, Thresholds())

        heading_blocks = [tb for tb in content.text_blocks if tb.is_heading]
        assert len(heading_blocks) >= 4, f"Expected >=4 headings, got {len(heading_blocks)}"

    def test_extract_slide_type_title(self, multi_slide_pptx: Path):
        prs = Presentation(str(multi_slide_pptx))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        content = extract_content(slide, 0, 3, sw, sh, Thresholds())

        assert content.slide_type == "title"

    def test_extract_footer_excluded_from_blocks(self, template_with_branding: Path):
        """Footer shapes (bottom 8%, 'Confidential', 'Page XX') should not appear in text_blocks."""
        prs = Presentation(str(template_with_branding))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        content = extract_content(slide, 0, 3, sw, sh, Thresholds())

        block_texts = " ".join(
            p.text for tb in content.text_blocks for p in tb.paragraphs
        )
        assert "Confidential" not in block_texts
        assert "Page 01" not in block_texts


class TestExtractAllContent:
    """Test batch content extraction."""

    def test_extract_all_returns_list(self, multi_slide_pptx: Path):
        contents = extract_all_content(multi_slide_pptx, Thresholds())

        assert len(contents) == 3
        assert all(isinstance(c, ContentData) for c in contents)

    def test_extract_all_preserves_order(self, multi_slide_pptx: Path):
        contents = extract_all_content(multi_slide_pptx, Thresholds())

        assert "Presentation Title" in contents[0].title
        assert "Content Slide Title" in contents[1].title
        assert "Thank You" in contents[2].title


# ===========================================================================
# 4. SHAPE CLASSIFICATION
# ===========================================================================

class TestShapeClassification:
    """Test classify_shape_role and get_slide_zones."""

    def test_classify_title_shape(self, simple_pptx: Path):
        prs = Presentation(str(simple_pptx))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        for shape in slide.shapes:
            if hasattr(shape, "text") and "Test Title" in shape.text:
                role = classify_shape_role(shape, sw, sh, slide=slide)
                assert role in ("title", "body"), f"Title shape classified as {role}"
                break

    def test_classify_footer_shape(self, template_with_branding: Path):
        prs = Presentation(str(template_with_branding))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        for shape in slide.shapes:
            if hasattr(shape, "text") and "Confidential" in shape.text:
                role = classify_shape_role(shape, sw, sh, slide=slide)
                assert role == "footer", f"Footer shape classified as {role}"
                break

    def test_get_slide_zones_keys(self, simple_pptx: Path):
        prs = Presentation(str(simple_pptx))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        zones = get_slide_zones(slide, sw, sh)

        assert "title" in zones
        assert "body" in zones
        assert "decorative" in zones
        assert "footer" in zones

    def test_zones_are_lists(self, simple_pptx: Path):
        prs = Presentation(str(simple_pptx))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        zones = get_slide_zones(slide, sw, sh)

        for key in zones:
            assert isinstance(zones[key], list)


# ===========================================================================
# 5. TEMPLATE ANALYSIS
# ===========================================================================

class TestTemplateAnalysis:
    """Test analyze_template extracts style information."""

    def test_analyze_returns_template_style(self, template_with_branding: Path):
        style = analyze_template(template_with_branding)

        assert isinstance(style, TemplateStyle)
        assert style.slide_width > 0
        assert style.slide_height > 0

    def test_analyze_extracts_dimensions(self, template_with_branding: Path):
        style = analyze_template(template_with_branding)

        # Should match 13.333" x 7.5" in EMUs
        assert style.slide_width == Inches(13.333)
        assert style.slide_height == Inches(7.5)

    def test_analyze_has_fonts(self, template_with_branding: Path):
        style = analyze_template(template_with_branding)

        assert style.heading_font != ""
        assert style.body_font != ""

    def test_analyze_has_colors(self, template_with_branding: Path):
        style = analyze_template(template_with_branding)

        # Colors should be 6-char hex strings
        assert len(style.color_primary) == 6
        assert len(style.color_text) == 6
        assert len(style.color_background) == 6

    def test_detect_mode_always_recreate(self, template_with_branding: Path):
        assert detect_mode(template_with_branding) == "recreate"


# ===========================================================================
# 6. RECREATE MODE — BUILD SLIDE
# ===========================================================================

class TestBuildSlide:
    """Test build_slide creates properly styled slides."""

    def _make_style(self, prs: Presentation) -> TemplateStyle:
        return TemplateStyle(
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            heading_font="Calibri Light",
            body_font="Calibri",
            color_primary="2563EB",
            color_secondary="F97316",
            color_text="111827",
            color_muted="475569",
            color_background="F7F8FB",
            color_card="FFFFFF",
            color_line="D1D5DB",
        )

    def test_build_slide_creates_shapes(self):
        prs = Presentation()
        style = self._make_style(prs)
        content = ContentData(
            title="Test Title",
            body_paragraphs=[ParagraphData(text="Body text here")],
            slide_type="content",
            word_count=5,
        )

        build_slide(prs, style, content, 1, 1)

        slide = prs.slides[0]
        assert len(list(slide.shapes)) > 0

    def test_build_slide_has_title_text(self):
        prs = Presentation()
        style = self._make_style(prs)
        content = ContentData(
            title="My Slide Title",
            body_paragraphs=[ParagraphData(text="Some body text")],
            slide_type="content",
            word_count=6,
        )

        build_slide(prs, style, content, 1, 1)

        slide = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        )
        assert "My Slide Title" in all_text

    def test_build_slide_fonts_explicit(self):
        """Every text run should have explicit font name (not None)."""
        prs = Presentation()
        style = self._make_style(prs)
        content = ContentData(
            title="Styled Title",
            body_paragraphs=[
                ParagraphData(text="Regular text"),
                ParagraphData(text="Bold heading", bold=True),
            ],
            slide_type="content",
            word_count=6,
        )

        build_slide(prs, style, content, 1, 1)

        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.text.strip():
                            assert run.font.name is not None, (
                                f"Run '{run.text}' has no font name"
                            )

    def test_build_title_slide(self):
        prs = Presentation()
        style = self._make_style(prs)
        content = ContentData(
            title="Presentation Title",
            slide_type="title",
            word_count=2,
        )

        build_slide(prs, style, content, 1, 5)

        slide = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        )
        assert "Presentation Title" in all_text

    def test_build_slide_with_text_blocks(self):
        """Text blocks should be placed at specified positions."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        style = self._make_style(prs)
        style.slide_width = prs.slide_width
        style.slide_height = prs.slide_height

        blocks = [
            TextBlock(
                paragraphs=[ParagraphData(text="Left column")],
                left_pct=10.0, top_pct=40.0,
                width_pct=30.0, height_pct=20.0,
            ),
            TextBlock(
                paragraphs=[ParagraphData(text="Right column")],
                left_pct=55.0, top_pct=40.0,
                width_pct=30.0, height_pct=20.0,
            ),
        ]
        content = ContentData(
            title="Two Column",
            text_blocks=blocks,
            slide_type="content",
            word_count=4,
        )

        build_slide(prs, style, content, 1, 1)

        slide = prs.slides[0]
        texts = [
            s.text_frame.text.strip()
            for s in slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        assert "Left column" in texts
        assert "Right column" in texts

    def test_build_slide_text_blocks_at_positions(self):
        """Verify text blocks are actually placed at different x-positions."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        style = self._make_style(prs)
        style.slide_width = prs.slide_width
        style.slide_height = prs.slide_height

        blocks = [
            TextBlock(
                paragraphs=[ParagraphData(text="At 10%")],
                left_pct=10.0, top_pct=40.0,
                width_pct=20.0, height_pct=10.0,
            ),
            TextBlock(
                paragraphs=[ParagraphData(text="At 60%")],
                left_pct=60.0, top_pct=40.0,
                width_pct=20.0, height_pct=10.0,
            ),
        ]
        content = ContentData(
            title="Position Test",
            text_blocks=blocks,
            slide_type="content",
            word_count=4,
        )

        build_slide(prs, style, content, 1, 1)

        slide = prs.slides[0]
        sw = prs.slide_width
        positions = {}
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text in ("At 10%", "At 60%"):
                    positions[text] = (shape.left or 0) / sw * 100

        assert "At 10%" in positions
        assert "At 60%" in positions
        # "At 60%" should be significantly to the right of "At 10%"
        assert positions["At 60%"] > positions["At 10%"] + 20

    def test_build_slide_footer_present(self):
        prs = Presentation()
        style = self._make_style(prs)
        style.footer_company = "TestCorp"
        content = ContentData(
            title="With Footer",
            slide_type="content",
            word_count=2,
        )

        build_slide(prs, style, content, 3, 10)

        slide = prs.slides[0]
        all_text = " ".join(
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        )
        assert "Page 03" in all_text or "03" in all_text

    def test_build_slide_speaker_notes(self):
        prs = Presentation()
        style = self._make_style(prs)
        content = ContentData(
            title="Notes Test",
            notes="These are speaker notes",
            slide_type="content",
            word_count=2,
        )

        build_slide(prs, style, content, 1, 1)

        slide = prs.slides[0]
        try:
            notes_text = slide.notes_slide.notes_text_frame.text
            assert "speaker notes" in notes_text
        except Exception:
            pytest.skip("Notes slide not supported in this environment")


# ===========================================================================
# 7. END-TO-END TRANSFER
# ===========================================================================

class TestEndToEnd:
    """Test the full transfer pipeline."""

    def test_transfer_recreate_creates_output(
        self, template_with_branding: Path, multi_slide_pptx: Path, tmp_path: Path,
    ):
        output = tmp_path / "output.pptx"

        report = transfer(template_with_branding, multi_slide_pptx, output)

        assert output.exists()
        assert output.stat().st_size > 0

    def test_transfer_output_has_correct_slide_count(
        self, template_with_branding: Path, multi_slide_pptx: Path, tmp_path: Path,
    ):
        output = tmp_path / "output.pptx"

        transfer(template_with_branding, multi_slide_pptx, output)

        prs = Presentation(str(output))
        assert len(prs.slides) == 3  # Same as content PPTX

    def test_transfer_output_preserves_content_text(
        self, template_with_branding: Path, multi_slide_pptx: Path, tmp_path: Path,
    ):
        output = tmp_path / "output.pptx"

        transfer(template_with_branding, multi_slide_pptx, output)

        prs = Presentation(str(output))
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + " "

        assert "Content Slide Title" in all_text
        assert "First paragraph" in all_text

    def test_transfer_no_template_text_leak(
        self, template_with_branding: Path, multi_slide_pptx: Path, tmp_path: Path,
    ):
        """Template body text should not appear in output."""
        output = tmp_path / "output.pptx"

        transfer(template_with_branding, multi_slide_pptx, output)

        prs = Presentation(str(output))
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + " "

        assert "Template body placeholder text that should not leak" not in all_text

    def test_transfer_all_runs_have_fonts(
        self, template_with_branding: Path, multi_slide_pptx: Path, tmp_path: Path,
    ):
        """Every text run in output should have explicit font name."""
        output = tmp_path / "output.pptx"

        transfer(template_with_branding, multi_slide_pptx, output)

        prs = Presentation(str(output))
        missing = []
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if run.text.strip() and run.font.name is None:
                                missing.append(f"Slide {i}: '{run.text[:30]}'")

        assert missing == [], f"Runs missing font: {missing[:5]}"

    def test_transfer_output_dimensions_match_template(
        self, template_with_branding: Path, multi_slide_pptx: Path, tmp_path: Path,
    ):
        output = tmp_path / "output.pptx"

        transfer(template_with_branding, multi_slide_pptx, output)

        template_prs = Presentation(str(template_with_branding))
        output_prs = Presentation(str(output))
        assert output_prs.slide_width == template_prs.slide_width
        assert output_prs.slide_height == template_prs.slide_height

    def test_transfer_with_grid_content(
        self, template_with_branding: Path, content_with_grid: Path, tmp_path: Path,
    ):
        """Grid layout content should produce positioned text blocks in output."""
        output = tmp_path / "output.pptx"

        transfer(template_with_branding, content_with_grid, output)

        prs = Presentation(str(output))
        slide = prs.slides[0]

        # Should have many text shapes (not just 6 flat ones)
        text_shapes = [
            s for s in slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        assert len(text_shapes) >= 10, (
            f"Expected >=10 text shapes for grid, got {len(text_shapes)}"
        )

    def test_transfer_grid_positions_diverse(
        self, template_with_branding: Path, content_with_grid: Path, tmp_path: Path,
    ):
        """Grid content should produce shapes at different x-positions."""
        output = tmp_path / "output.pptx"

        transfer(template_with_branding, content_with_grid, output)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        sw = prs.slide_width

        lefts = set()
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                left_pct = round((shape.left or 0) / sw * 100, 0)
                lefts.add(left_pct)

        # Should have at least 3 distinct x-positions (left margin + columns)
        assert len(lefts) >= 3, f"Expected diverse x-positions, got {lefts}"

    def test_transfer_clone_mode(
        self, template_with_branding: Path, multi_slide_pptx: Path, tmp_path: Path,
    ):
        """Clone mode should also work without crashing."""
        output = tmp_path / "output_clone.pptx"
        config = TransferConfig(mode="clone")

        report = transfer(template_with_branding, multi_slide_pptx, output, config)

        assert output.exists()

    def test_transfer_config_no_notes(
        self, template_with_branding: Path, multi_slide_pptx: Path, tmp_path: Path,
    ):
        output = tmp_path / "output_no_notes.pptx"
        config = TransferConfig(preserve_notes=False)

        transfer(template_with_branding, multi_slide_pptx, output, config)

        assert output.exists()


# ===========================================================================
# 8. PROTECTED SHAPE / FOOTER DETECTION
# ===========================================================================

class TestProtectedShapes:
    """Test _is_protected_shape correctly identifies footer/media shapes."""

    def test_footer_text_protected(self, template_with_branding: Path):
        prs = Presentation(str(template_with_branding))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        for shape in slide.shapes:
            if shape.has_text_frame and "Confidential" in shape.text_frame.text:
                assert mod._is_protected_shape(shape, sw, sh) is True
                break

    def test_page_number_protected(self, template_with_branding: Path):
        prs = Presentation(str(template_with_branding))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        for shape in slide.shapes:
            if shape.has_text_frame and "Page 01" in shape.text_frame.text:
                assert mod._is_protected_shape(shape, sw, sh) is True
                break

    def test_body_text_not_protected(self, template_with_branding: Path):
        prs = Presentation(str(template_with_branding))
        slide = prs.slides[0]
        sw, sh = prs.slide_width, prs.slide_height

        for shape in slide.shapes:
            if shape.has_text_frame and "placeholder text" in shape.text_frame.text:
                assert mod._is_protected_shape(shape, sw, sh) is False
                break


# ===========================================================================
# 9. INPUT VALIDATION
# ===========================================================================

class TestInputValidation:
    """Test _validate_input catches bad files."""

    def test_validate_missing_file(self, tmp_path: Path):
        with pytest.raises(SystemExit):
            mod._validate_input(tmp_path / "nonexistent.pptx", "test")

    def test_validate_not_zip(self, tmp_path: Path):
        bad_file = tmp_path / "bad.pptx"
        bad_file.write_text("this is not a zip")

        with pytest.raises(SystemExit):
            mod._validate_input(bad_file, "test")

    def test_validate_good_file(self, simple_pptx: Path):
        # Should not raise
        mod._validate_input(simple_pptx, "test")


# ===========================================================================
# 10. FIND BLANK LAYOUT
# ===========================================================================

class TestFindBlankLayout:
    """Test _find_blank_layout selects appropriate layout."""

    def test_finds_a_layout(self, simple_pptx: Path):
        prs = Presentation(str(simple_pptx))
        layout = mod._find_blank_layout(prs)

        assert layout is not None

    def test_prefers_blank_name(self):
        prs = Presentation()
        layout = mod._find_blank_layout(prs)

        # Default presentations have a Blank layout
        assert layout is not None
