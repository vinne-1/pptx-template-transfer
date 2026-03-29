"""Regression tests using real PPTX decks.

These tests verify that the transfer pipeline works correctly for any
source+target pair by checking:
  - Source content is the primary authority in the output
  - No target body-text contamination
  - No text leakage across slides
  - Target shell (fonts, colors, footer) is preserved
  - Semantic slide-type classification works
  - Source coverage is adequate
  - Content provenance is tracked

The specific decks used here are a *regression case*, not hardcoded
assumptions — the assertions are generic and should hold for ANY
valid source+target pair.

Parameter convention:
  template = STYLE source (target visual style deck)
  content  = CONTENT source (source content deck)
"""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

# ---------------------------------------------------------------------------
# Deck paths — skip entire module if not present
# ---------------------------------------------------------------------------
REPO_ROOT = Path(__file__).resolve().parent.parent

# SOURCE deck = Aaritya MDR (16 slides of MDR report content)
SOURCE_PATH = REPO_ROOT / "Aaritya_MDR.pptx"
# TARGET deck = Airowire (7 slides with nice visual style)
TARGET_PATH = REPO_ROOT / "Airowire_Day2_Security_Flow_reporting_update.pptx"

HAVE_DECKS = SOURCE_PATH.exists() and TARGET_PATH.exists()
pytestmark = pytest.mark.skipif(not HAVE_DECKS, reason="Regression decks not found")

from pptx_template_transfer import (
    TransferConfig,
    Thresholds,
    extract_all_content,
    transfer,
)
from pptx_template_transfer.analysis.theme_extractor import (
    analyze_template,
    _extract_footer_text,
)
from pptx_template_transfer.validation.quality_report import (
    generate_quality_report,
    _detect_text_leakage,
    _detect_body_in_forbidden_zones,
)
from pptx_template_transfer.validation.contamination_checker import (
    check_target_contamination,
)
from pptx_template_transfer.validation.source_coverage import (
    compute_source_coverage,
)
from pptx_template_transfer.helpers import text_of, word_count


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(scope="module")
def output_path(tmp_path_factory) -> Path:
    """Run the transfer once with CORRECT parameter order and return output."""
    out = tmp_path_factory.mktemp("regression") / "output.pptx"
    # template = STYLE source (Airowire), content = CONTENT source (Aaritya)
    config = TransferConfig(mode="recreate")
    transfer(TARGET_PATH, SOURCE_PATH, out, config)
    return out


@pytest.fixture(scope="module")
def output_prs(output_path: Path) -> Presentation:
    return Presentation(str(output_path))


@pytest.fixture(scope="module")
def source_content() -> list:
    return extract_all_content(SOURCE_PATH, Thresholds())


@pytest.fixture(scope="module")
def template_style():
    return analyze_template(TARGET_PATH)


@pytest.fixture(scope="module")
def quality_report(output_prs, source_content):
    return generate_quality_report(output_prs, source_content)


# ---------------------------------------------------------------------------
# A. Source content is primary authority
# ---------------------------------------------------------------------------

class TestSourceContentAuthority:
    """Output must contain source (Aaritya MDR) content, not target (Airowire)."""

    def test_output_slide_count_matches_source(
        self, output_prs: Presentation, source_content,
    ):
        """Output should have same number of slides as source deck."""
        assert len(output_prs.slides) == len(source_content), (
            f"Output has {len(output_prs.slides)} slides but source has "
            f"{len(source_content)} — content may come from wrong deck"
        )

    def test_source_titles_in_output(self, output_prs: Presentation, source_content):
        """Source slide titles should appear in the output."""
        output_texts: list[str] = []
        for slide in output_prs.slides:
            for shape in slide.shapes:
                t = text_of(shape)
                if t:
                    output_texts.append(t.lower())
        all_output = " ".join(output_texts)

        found = 0
        for cd in source_content:
            if cd.title and cd.title.lower()[:30] in all_output:
                found += 1

        coverage = found / len(source_content) if source_content else 0
        assert coverage >= 0.6, (
            f"Only {found}/{len(source_content)} source titles found in output "
            f"({coverage:.0%}) — output may be dominated by target content"
        )

    def test_source_body_text_in_output(self, output_prs: Presentation, source_content):
        """Key source body phrases should appear in output."""
        all_output = ""
        for slide in output_prs.slides:
            for shape in slide.shapes:
                t = text_of(shape)
                if t:
                    all_output += " " + t.lower()

        # Sample distinctive source phrases (from Aaritya MDR content)
        # These are generic enough to not be hardcoded to specific content
        source_phrases = set()
        for cd in source_content:
            for p in cd.body_paragraphs[:3]:
                words = p.text.split()
                if len(words) >= 4:
                    source_phrases.add(" ".join(words[:4]).lower())

        found = sum(1 for phrase in source_phrases if phrase in all_output)
        if source_phrases:
            coverage = found / len(source_phrases)
            assert coverage >= 0.3, (
                f"Only {found}/{len(source_phrases)} source phrases found "
                f"({coverage:.0%})"
            )


# ---------------------------------------------------------------------------
# B. No target-content contamination
# ---------------------------------------------------------------------------

class TestNoTargetContamination:
    """Output body text must not come from the target/style deck."""

    def test_no_target_contamination(self, output_prs: Presentation):
        """Jaccard similarity between output and target body text must be low."""
        warnings = check_target_contamination(output_prs, TARGET_PATH)
        assert warnings == [], (
            f"Target-content contamination detected:\n" +
            "\n".join(f"  {w}" for w in warnings)
        )

    def test_no_footer_text_leak(self, output_prs: Presentation):
        """Source body content must not be promoted to template footer."""
        for i, slide in enumerate(output_prs.slides):
            for shape in slide.shapes:
                t = text_of(shape)
                if "Complete coverage is necessary for maximum protection" in t:
                    pytest.fail(
                        f"Source body text leaked to slide {i+1} as footer"
                    )


# ---------------------------------------------------------------------------
# C. Source coverage
# ---------------------------------------------------------------------------

class TestSourceCoverage:
    """Adequate source content must make it into the output."""

    def test_overall_source_coverage(self, output_prs: Presentation, source_content):
        """Overall text coverage of source content must be >= 50%."""
        cov = compute_source_coverage(output_prs, source_content)
        assert cov.overall_text_coverage_pct >= 50, (
            f"Source coverage {cov.overall_text_coverage_pct}% is below 50%"
        )

    def test_no_unmapped_source_slides(self, output_prs: Presentation, source_content):
        """Every source slide should map to at least one output slide."""
        cov = compute_source_coverage(output_prs, source_content)
        assert len(cov.unmapped_source_slides) == 0, (
            f"Unmapped source slides: {[s+1 for s in cov.unmapped_source_slides]}"
        )

    def test_provenance_in_report(self, output_path: Path):
        """The transfer report should include provenance and coverage data."""
        config = TransferConfig(mode="recreate")
        report = transfer(TARGET_PATH, SOURCE_PATH, output_path, config)
        assert "source_coverage" in report
        assert report["source_coverage"]["overall_pct"] > 0
        for slide_report in report["slides"]:
            assert "provenance" in slide_report
            assert slide_report["provenance"]["body"] == "source_content"
            assert slide_report["provenance"]["footer"] == "target_shell"


# ---------------------------------------------------------------------------
# D. Target shell preserved
# ---------------------------------------------------------------------------

class TestTargetShellPreserved:
    """Output must reflect the target template's visual DNA."""

    def test_output_uses_template_dimensions(
        self, output_prs: Presentation, template_style,
    ):
        assert output_prs.slide_width == template_style.slide_width
        assert output_prs.slide_height == template_style.slide_height

    def test_template_fonts_detected(self, template_style):
        assert template_style.heading_font, "No heading font detected"
        assert template_style.body_font, "No body font detected"

    def test_target_footer_in_output(self, output_prs: Presentation, template_style):
        """Target footer (shell element) should appear in output."""
        if not template_style.footer_company:
            pytest.skip("No footer detected in target template")
        found = False
        for slide in output_prs.slides:
            for shape in slide.shapes:
                if template_style.footer_company in text_of(shape):
                    found = True
                    break
            if found:
                break
        assert found, f"Target footer '{template_style.footer_company}' not in output"


# ---------------------------------------------------------------------------
# E. Slide-type classification
# ---------------------------------------------------------------------------

class TestSlideTypeClassification:

    def test_all_slides_have_types(self, source_content):
        for i, cd in enumerate(source_content):
            assert cd.slide_type, f"Slide {i+1} has no slide_type"

    def test_first_slide_is_title_type(self, source_content):
        assert source_content[0].slide_type == "title"

    def test_slide_types_are_diverse(self, source_content):
        types = set(cd.slide_type for cd in source_content)
        assert len(types) >= 2


# ---------------------------------------------------------------------------
# F. Quality validation
# ---------------------------------------------------------------------------

class TestQualityValidation:

    def test_quality_score_above_minimum(self, quality_report):
        assert quality_report.overall_score >= 50, (
            f"Quality score {quality_report.overall_score}/100 below minimum"
        )

    def test_all_slides_built_natively(self, quality_report):
        assert quality_report.fallback_count == 0

    def test_no_empty_slides(self, output_prs: Presentation):
        for i, slide in enumerate(output_prs.slides):
            texts = [text_of(s) for s in slide.shapes if text_of(s)]
            assert texts, f"Slide {i+1} has no text content"

    def test_no_duplicate_leakage(self, output_prs: Presentation, source_content):
        warnings = _detect_text_leakage(output_prs, source_content)
        assert warnings == [], "\n".join(warnings)


# ---------------------------------------------------------------------------
# G. Section label quality
# ---------------------------------------------------------------------------

class TestSectionLabels:
    """Section labels must be clean, professional, and derived from source."""

    def test_no_broken_labels(self, output_prs: Presentation):
        """No label should end with ?, !, or trailing prepositions."""
        import re
        bad_endings = re.compile(r"[?!:;,.\-–—]\s*$")
        trailing_junk = {"the", "a", "an", "of", "for", "in", "on", "to", "and", "with", "by"}
        for i, slide in enumerate(output_prs.slides):
            for shape in slide.shapes:
                t = text_of(shape).strip()
                if not t or not t.isupper() or len(t) > 40:
                    continue
                if t in ("AIROWIRE SECURITY SERVICES",):
                    continue
                assert not bad_endings.search(t), (
                    f"Slide {i+1}: broken label '{t}' ends with punctuation"
                )
                words = t.split()
                if words:
                    assert words[-1].lower() not in trailing_junk, (
                        f"Slide {i+1}: label '{t}' ends with junk word"
                    )

    def test_incident_slides_labeled_incident(self, output_prs: Presentation, source_content):
        """Slides with 'Incident' in title should have INCIDENT label."""
        for i, cd in enumerate(source_content):
            if "incident" in (cd.title or "").lower():
                slide = output_prs.slides[i]
                labels = [
                    text_of(s).strip() for s in slide.shapes
                    if text_of(s).strip().isupper() and len(text_of(s).strip()) < 30
                    and text_of(s).strip() not in ("AIROWIRE SECURITY SERVICES",)
                ]
                assert any("INCIDENT" in l for l in labels), (
                    f"Slide {i+1} ({cd.title[:40]}) should have INCIDENT label, "
                    f"got {labels}"
                )


# ---------------------------------------------------------------------------
# H. Quality scoring
# ---------------------------------------------------------------------------

class TestQualityScoring:

    def test_quality_report_has_per_slide_scores(self, quality_report):
        """Every slide in quality report should have a composite score."""
        for sq in quality_report.slides:
            score = getattr(sq, "slide_score", None)
            assert score is not None, f"Slide {sq.slide_index+1} missing slide_score"
            assert 0 <= score <= 100

    def test_transfer_report_includes_quality(self, output_path: Path):
        """Transfer report should include quality section."""
        config = TransferConfig(mode="recreate")
        report = transfer(TARGET_PATH, SOURCE_PATH, output_path, config)
        assert "quality" in report
        assert "overall_score" in report["quality"]
        assert report["quality"]["overall_score"] > 0

    def test_no_slides_below_acceptance_gate(self, quality_report):
        """No slide should have quality score below 40 (hard fail)."""
        for sq in quality_report.slides:
            score = getattr(sq, "slide_score", sq.content_coverage_pct)
            assert score >= 40, (
                f"Slide {sq.slide_index+1} scored {score:.0f}/100, below hard floor"
            )


# ---------------------------------------------------------------------------
# I. Generalization tests
# ---------------------------------------------------------------------------

class TestGeneralization:

    def test_footer_detection_requires_multi_slide(self):
        """Single-slide bottom text is not a footer."""
        from pptx.util import Inches
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        sh = prs.slide_height
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(
            Inches(0.5), int(sh * 0.92), Inches(5), Inches(0.3),
        )
        tb.text_frame.text = "This sentence appears on only one slide."
        assert _extract_footer_text(prs) == ""

    def test_footer_detection_accepts_multi_slide(self):
        """Text on 2+ slides in footer zone IS a footer."""
        from pptx.util import Inches
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        sh = prs.slide_height
        for _ in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tb = slide.shapes.add_textbox(
                Inches(0.5), int(sh * 0.92), Inches(5), Inches(0.3),
            )
            tb.text_frame.text = "Acme Corp Confidential Footer"
        assert _extract_footer_text(prs) == "Acme Corp Confidential Footer"

    def test_extract_handles_empty_slide(self):
        from pptx_template_transfer import extract_content
        prs = Presentation()
        prs.slide_width = 9144000
        prs.slide_height = 6858000
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cd = extract_content(slide, 0, 1, prs.slide_width, prs.slide_height, Thresholds())
        assert cd.word_count == 0


# ---------------------------------------------------------------------------
# J. Reverse direction regression (Airowire content → Aaritya style)
# ---------------------------------------------------------------------------

class TestReverseTransfer:
    """Transfer in reverse direction must also produce clean output."""

    @pytest.fixture(scope="class")
    def reverse_result(self, tmp_path_factory):
        out = tmp_path_factory.mktemp("reverse") / "output_reverse.pptx"
        # Aaritya as STYLE source, Airowire as CONTENT source
        config = TransferConfig(mode="recreate")
        return transfer(SOURCE_PATH, TARGET_PATH, out, config)

    @pytest.fixture(scope="class")
    def reverse_prs(self, reverse_result, tmp_path_factory):
        # Re-read the generated file
        out = tmp_path_factory.mktemp("reverse_read") / "output_reverse.pptx"
        config = TransferConfig(mode="recreate")
        transfer(SOURCE_PATH, TARGET_PATH, out, config)
        return Presentation(str(out))

    def test_reverse_slide_count(self, reverse_result):
        """Reverse transfer should produce slides matching Airowire content count."""
        slides = reverse_result["slides"]
        assert len(slides) == 7  # Airowire has 7 slides

    def test_reverse_no_errors(self, reverse_result):
        """No slide build errors in reverse direction."""
        assert not reverse_result.get("errors", [])

    def test_reverse_quality_above_threshold(self, reverse_result):
        """Quality score should be above 70 even in reverse direction."""
        assert reverse_result["quality"]["overall_score"] >= 70

    def test_reverse_coverage_above_threshold(self, reverse_result):
        """Source coverage should be above 80%."""
        assert reverse_result["source_coverage"]["overall_pct"] >= 80

    def test_reverse_no_dropped_content(self, reverse_result):
        """No false 'dropped content' messages (images/tables are transferred)."""
        for s in reverse_result["slides"]:
            dropped = s.get("dropped_content", [])
            # Only charts should appear as dropped
            for d in dropped:
                assert "chart" in d.lower(), f"Unexpected dropped content: {d}"


# ---------------------------------------------------------------------------
# K. Edge-case robustness
# ---------------------------------------------------------------------------

class TestEdgeCases:
    """Stress tests for edge-case inputs."""

    def test_single_slide_source(self, tmp_path):
        """Transfer with a 1-slide source should not crash."""
        from pptx.util import Inches
        # Create a minimal 1-slide source
        src = Presentation()
        src.slide_width = Inches(13.333)
        src.slide_height = Inches(7.5)
        slide = src.slides.add_slide(src.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tb.text_frame.text = "Single Slide Content"
        src_path = tmp_path / "single.pptx"
        src.save(str(src_path))

        out = tmp_path / "out_single.pptx"
        result = transfer(TARGET_PATH, src_path, out)
        assert len(result["slides"]) == 1
        assert not result.get("errors", [])

    def test_empty_body_slide(self, tmp_path):
        """A slide with title but no body should render without error."""
        from pptx.util import Inches, Pt as PxPt
        src = Presentation()
        src.slide_width = Inches(13.333)
        src.slide_height = Inches(7.5)
        slide = src.slides.add_slide(src.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = "Title Only Slide"
        p.font.size = PxPt(28)
        src_path = tmp_path / "title_only.pptx"
        src.save(str(src_path))

        out = tmp_path / "out_title_only.pptx"
        result = transfer(TARGET_PATH, src_path, out)
        assert len(result["slides"]) == 1
        assert not result.get("errors", [])

    def test_very_long_text(self, tmp_path):
        """A slide with very long body text should not crash or overflow."""
        from pptx.util import Inches
        src = Presentation()
        src.slide_width = Inches(13.333)
        src.slide_height = Inches(7.5)
        slide = src.slides.add_slide(src.slide_layouts[6])
        # Title
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
        tb.text_frame.text = "Stress Test Slide"
        # Long body
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
        tb2.text_frame.word_wrap = True
        long_text = "This is a test sentence with enough words to matter. " * 50
        tb2.text_frame.text = long_text
        src_path = tmp_path / "long_text.pptx"
        src.save(str(src_path))

        out = tmp_path / "out_long.pptx"
        result = transfer(TARGET_PATH, src_path, out)
        assert not result.get("errors", [])
        # Quality should still be reasonable
        assert result["quality"]["overall_score"] >= 50

    def test_many_slides(self, tmp_path):
        """A source with 30+ slides should transfer without error."""
        from pptx.util import Inches
        src = Presentation()
        src.slide_width = Inches(13.333)
        src.slide_height = Inches(7.5)
        for i in range(30):
            slide = src.slides.add_slide(src.slide_layouts[6])
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            tb.text_frame.text = f"Slide {i+1} content paragraph about topic {i+1}."
        src_path = tmp_path / "many_slides.pptx"
        src.save(str(src_path))

        out = tmp_path / "out_many.pptx"
        result = transfer(TARGET_PATH, src_path, out)
        assert len(result["slides"]) == 30
        assert not result.get("errors", [])

    def test_special_characters(self, tmp_path):
        """Unicode and special chars in content should not crash."""
        from pptx.util import Inches
        src = Presentation()
        src.slide_width = Inches(13.333)
        src.slide_height = Inches(7.5)
        slide = src.slides.add_slide(src.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tb.text_frame.text = "Ünïcödé: 日本語テスト • bullets — dashes « quotes » ñ ø å"
        src_path = tmp_path / "unicode.pptx"
        src.save(str(src_path))

        out = tmp_path / "out_unicode.pptx"
        result = transfer(TARGET_PATH, src_path, out)
        assert not result.get("errors", [])
